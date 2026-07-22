#!/usr/bin/env python3
"""Shared 'Good Code, Good Vibes' deck theme + templates. Import with:
    from deck_common import *
    prs = init_deck()
    ... build slides ...
    save('/path/out.pptx')
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml import parse_xml
import os, re, subprocess

ASSETS = os.path.join(os.path.dirname(__file__), "assets")

# ---- palette ----
BG    = RGBColor(0x0B, 0x0E, 0x14)
PANEL = RGBColor(0x15, 0x1B, 0x26)
GREEN = RGBColor(0x00, 0xFF, 0x41)   # good / defense
RED   = RGBColor(0xFF, 0x4D, 0x4D)   # harm / problem
WHITE = RGBColor(0xF5, 0xF7, 0xFA)
MUTED = RGBColor(0x8A, 0x94, 0xA6)
DIM   = RGBColor(0x5A, 0x64, 0x74)
F_HEAD = "Arial"; F_BODY = "Arial"; F_MONO = "Courier New"

prs = None; SW = None; SH = None; BLANK = None

def init_deck():
    global prs, SW, SH, BLANK
    prs = Presentation()
    prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
    SW, SH = prs.slide_width, prs.slide_height
    BLANK = prs.slide_layouts[6]
    return prs

def save(path):
    os.makedirs(os.path.dirname(path), exist_ok=True)
    prs.save(path)
    return len(prs.slides._sldIdLst)

# ---- primitives ----
def slide(bg=BG):
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid(); s.background.fill.fore_color.rgb = bg
    return s

def _set_font(run, size, color, bold, font):
    run.font.size = Pt(size); run.font.bold = bold
    run.font.color.rgb = color; run.font.name = font

def text(s, l, t, w, h, runs, size=18, color=WHITE, bold=False, font=F_BODY,
         align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, spacing=1.0):
    tb = s.shapes.add_textbox(l, t, w, h); tf = tb.text_frame
    tf.word_wrap = True; tf.vertical_anchor = anchor
    if isinstance(runs, str): runs = [runs]
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.line_spacing = spacing
        segs = para if isinstance(para, list) else [(para, {})]
        for seg_txt, opts in segs:
            r = p.add_run(); r.text = seg_txt
            _set_font(r, opts.get("size", size), opts.get("color", color),
                      opts.get("bold", bold), opts.get("font", font))
    return tb

def rect(s, l, t, w, h, fill=PANEL, line=None, line_w=1.0, rounded=False, dash=False):
    shp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE, l, t, w, h)
    if fill is None: shp.fill.background()
    else: shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None: shp.line.fill.background()
    else:
        shp.line.color.rgb = line; shp.line.width = Pt(line_w)
        if dash:
            shp.line._get_or_add_ln().append(parse_xml(
                '<a:prstDash xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" val="dash"/>'))
    shp.shadow.inherit = False
    return shp

def kicker(s, txt, color=GREEN):
    text(s, Inches(0.7), Inches(0.55), Inches(11), Inches(0.4),
         [[(txt, {})]], size=13, color=color, bold=True, font=F_MONO)

def footer(s, txt):
    text(s, Inches(0.7), Inches(6.95), Inches(12), Inches(0.4),
         [[(txt, {})]], size=10, color=DIM, font=F_MONO)

def accent(s, l=Inches(0.7), t=Inches(1.02), w=Inches(0.9), h=Inches(0.08), color=GREEN):
    rect(s, l, t, w, h, fill=color)

def section(kick, title, color=GREEN, sub=None):
    """Act divider — a full-bleed section break, color-coded (red=problem, green=defense)."""
    s = slide(); rect(s, 0, 0, Inches(0.28), SH, fill=color)
    text(s, Inches(0.9), Inches(2.4), Inches(11.5), Inches(0.6),
         [[(kick, {"size": 18, "color": color, "bold": True, "font": F_MONO})]])
    text(s, Inches(0.88), Inches(3.0), Inches(11.5), Inches(1.7),
         [[(title, {"size": 48, "bold": True, "color": WHITE})]])
    if sub:
        text(s, Inches(0.92), Inches(4.7), Inches(11), Inches(0.7),
             [[(sub, {"size": 19, "color": MUTED})]])
    return s

def notes(s, txt):
    s.notes_slide.notes_text_frame.text = txt

# ---- templates ----
def cover(title_lines, subtitle, tag):
    s = slide(); rect(s, 0, 0, Inches(0.28), SH, fill=GREEN)
    text(s, Inches(0.9), Inches(2.2), Inches(11.5), Inches(2.2),
         [[(title_lines[0], {"size": 54, "bold": True, "color": WHITE})],
          [(title_lines[1], {"size": 54, "bold": True, "color": GREEN})]], spacing=1.02)
    text(s, Inches(0.92), Inches(4.5), Inches(11), Inches(0.6), [[(subtitle, {"size": 20, "color": MUTED})]])
    text(s, Inches(0.92), Inches(5.15), Inches(11), Inches(0.5), [[(tag, {"size": 13, "color": DIM, "font": F_MONO})]])
    return s

def divider(day, title, presenter=None):
    s = slide(); rect(s, 0, 0, Inches(0.28), SH, fill=GREEN)
    text(s, Inches(0.9), Inches(2.5), Inches(11.5), Inches(0.5),
         [[(day, {"size": 16, "color": GREEN, "bold": True, "font": F_MONO})]])
    text(s, Inches(0.88), Inches(3.05), Inches(11.5), Inches(1.6),
         [[(title, {"size": 46, "bold": True, "color": WHITE})]])
    if presenter:
        text(s, Inches(0.92), Inches(4.55), Inches(11), Inches(0.5),
             [[(presenter, {"size": 15, "color": MUTED})]])
    return s

def content(kick, title, body_paras, foot=None, title_size=34):
    s = slide(); kicker(s, kick); accent(s)
    text(s, Inches(0.7), Inches(1.3), Inches(12), Inches(1.2),
         [[(title, {"size": title_size, "bold": True, "color": WHITE})]])
    if body_paras:
        text(s, Inches(0.72), Inches(2.5), Inches(11.9), Inches(4.0), body_paras, size=20, color=WHITE, spacing=1.12)
    if foot: footer(s, foot)
    return s

def bullets(kick, title, items, foot=None):
    paras = [[("•  ", {"color": GREEN, "bold": True}), (it, {})] for it in items]
    return content(kick, title, paras, foot)

def big_question(kick, questions, foot=None):
    s = slide(); kicker(s, kick)
    text(s, Inches(1.2), Inches(0.6), Inches(3), Inches(2), [[("?", {"size": 120, "bold": True, "color": GREEN})]])
    paras = [[(q, {})] for q in questions]
    text(s, Inches(1.2), Inches(3.0), Inches(11), Inches(3.4), paras, size=30, color=WHITE, bold=True, spacing=1.12)
    if foot: footer(s, foot)
    return s

def media(kick, title, label, hint, foot=None):
    s = slide(); kicker(s, kick); accent(s)
    text(s, Inches(0.7), Inches(1.3), Inches(12), Inches(1.0), [[(title, {"size": 30, "bold": True, "color": WHITE})]])
    rect(s, Inches(2.4), Inches(2.55), Inches(8.5), Inches(3.9), fill=PANEL, line=GREEN, line_w=1.5, rounded=True, dash=True)
    text(s, Inches(2.4), Inches(3.7), Inches(8.5), Inches(0.8),
         [[(label, {"size": 22, "bold": True, "color": GREEN, "font": F_MONO})]], align=PP_ALIGN.CENTER)
    text(s, Inches(2.4), Inches(4.5), Inches(8.5), Inches(0.8), [[(hint, {"size": 14, "color": MUTED})]], align=PP_ALIGN.CENTER)
    if foot: footer(s, foot)
    return s

def columns(kick, title, cols, foot=None, accent_color=GREEN):
    s = slide(); kicker(s, kick, color=accent_color); accent(s, color=accent_color)
    text(s, Inches(0.7), Inches(1.3), Inches(12), Inches(1.0), [[(title, {"size": 32, "bold": True, "color": WHITE})]])
    n = len(cols); gap = Inches(0.4); left = Inches(0.7)
    total = SW - Inches(1.4) - gap * (n - 1); cw = Emu(int(total / n))
    top = Inches(2.7); ch = Inches(3.7)
    for i, (label, head, desc) in enumerate(cols):
        x = Emu(int(left + i * (cw + gap)))
        rect(s, x, top, cw, ch, fill=PANEL, rounded=True)
        rect(s, x, top, cw, Inches(0.09), fill=accent_color)
        text(s, Emu(int(x + Inches(0.3))), Emu(int(top + Inches(0.35))), Emu(int(cw - Inches(0.6))), Inches(0.5),
             [[(label.upper(), {"size": 12, "bold": True, "color": accent_color, "font": F_MONO})]])
        text(s, Emu(int(x + Inches(0.3))), Emu(int(top + Inches(0.85))), Emu(int(cw - Inches(0.6))), Inches(0.9),
             [[(head, {"size": 20, "bold": True, "color": WHITE})]])
        text(s, Emu(int(x + Inches(0.3))), Emu(int(top + Inches(1.8))), Emu(int(cw - Inches(0.6))), Inches(1.7),
             [[(desc, {"size": 15, "color": MUTED})]], spacing=1.1)
    if foot: footer(s, foot)
    return s

def checklist(kick, title, items, foot=None):
    s = slide(); kicker(s, kick); accent(s)
    text(s, Inches(0.7), Inches(1.3), Inches(12), Inches(1.0), [[(title, {"size": 34, "bold": True, "color": WHITE})]])
    paras = [[("✓  ", {"color": GREEN, "bold": True, "size": 20}), (it, {"size": 20})] for it in items]
    text(s, Inches(0.9), Inches(2.5), Inches(11.5), Inches(4.2), paras, size=20, color=WHITE, spacing=1.35)
    if foot: footer(s, foot)
    return s

def value_grid(kick, title, cards, cols=4, foot=None, accent_color=GREEN):
    """Taxonomy overview: N cards in a grid, all visible at once.
    cards = list of (icon, name, line). accent_color: RED for harms, GREEN for defenses."""
    import math
    s = slide(); kicker(s, kick, color=accent_color); accent(s, color=accent_color)
    text(s, Inches(0.7), Inches(1.12), Inches(12), Inches(0.9),
         [[(title, {"size": 30, "bold": True, "color": WHITE})]])
    rows = math.ceil(len(cards) / cols)
    gx, gy = Inches(0.3), Inches(0.3)
    gleft, gtop = Inches(0.7), Inches(2.2)
    gw = SW - Inches(1.4); gh = SH - gtop - Inches(0.7)
    cw = Emu(int((gw - gx * (cols - 1)) / cols))
    ch = Emu(int((gh - gy * (rows - 1)) / rows))
    for i, card in enumerate(cards):
        icon, name, harm = card
        r, c = divmod(i, cols)
        x = Emu(int(gleft + c * (cw + gx)))
        y = Emu(int(gtop + r * (ch + gy)))
        rect(s, x, y, cw, ch, fill=PANEL, rounded=True)
        rect(s, x, y, Inches(0.09), ch, fill=accent_color)  # left accent bar
        pad = Inches(0.28)
        text(s, Emu(int(x + pad)), Emu(int(y + Inches(0.24))), Emu(int(cw - pad * 2)), Inches(0.85),
             [[(f"{icon}  ", {"size": 19}), (name, {"size": 16, "bold": True, "color": WHITE})]], spacing=1.02)
        text(s, Emu(int(x + pad)), Emu(int(y + Inches(1.28))), Emu(int(cw - pad * 2)), Emu(int(ch - Inches(1.5))),
             [[(harm, {"size": 12.5, "color": MUTED})]], spacing=1.08)
    if foot: footer(s, foot)
    return s

def _img_dims(path):
    out = subprocess.run(["sips", "-g", "pixelWidth", "-g", "pixelHeight", path],
                         capture_output=True, text=True).stdout
    w = int(re.search(r"pixelWidth: (\d+)", out).group(1))
    h = int(re.search(r"pixelHeight: (\d+)", out).group(1))
    return w, h

def photo_fill(s, path, l, t, w, h):
    """Add an image cropped to completely fill the box (l,t,w,h) without distortion."""
    iw, ih = _img_dims(path)
    box_ar = w / h; img_ar = iw / ih
    pic = s.shapes.add_picture(path, l, t, width=w, height=h)
    if img_ar > box_ar:
        crop = (1 - box_ar / img_ar) / 2; pic.crop_left = crop; pic.crop_right = crop
    else:
        crop = (1 - img_ar / box_ar) / 2; pic.crop_top = crop; pic.crop_bottom = crop
    return pic

def photo_split(kick, title, body_paras, img_path, credit, foot=None, accent_color=GREEN):
    """Image fills the right ~42%; kicker/title/body on the left. Credit caption over the image."""
    s = slide(); kicker(s, kick, color=accent_color); accent(s, color=accent_color)
    half = Inches(5.6)
    photo_fill(s, os.path.join(ASSETS, img_path), Emu(int(SW - half)), 0, half, SH)
    tx = Inches(0.7); tw = Emu(int(SW - half - Inches(1.15)))
    text(s, tx, Inches(1.3), tw, Inches(1.4), [[(title, {"size": 30, "bold": True, "color": WHITE})]])
    if body_paras:
        text(s, tx, Inches(2.8), tw, Inches(3.4), body_paras, size=19, color=WHITE, spacing=1.16)
    rect(s, Emu(int(SW - half)), Emu(int(SH - Inches(0.34))), half, Inches(0.34), fill=BG)
    text(s, Emu(int(SW - half + Inches(0.15))), Emu(int(SH - Inches(0.32))), Emu(int(half - Inches(0.3))), Inches(0.3),
         [[(credit, {"size": 8.5, "color": MUTED, "font": F_MONO})]])
    if foot: footer(s, foot)
    return s

def headline_cards(kick, title, cards, foot=None, accent_color=RED):
    """A press wall: cards = (source, headline, value_tag). Reads like cited news clippings."""
    s = slide(); kicker(s, kick, color=accent_color); accent(s, color=accent_color)
    text(s, Inches(0.7), Inches(1.25), Inches(12), Inches(0.9), [[(title, {"size": 30, "bold": True, "color": WHITE})]])
    n = len(cards); gap = Inches(0.4); left = Inches(0.7)
    total = SW - Inches(1.4) - gap * (n - 1); cw = Emu(int(total / n))
    top = Inches(2.45); ch = Inches(4.0)
    for i, (source, headline, tag) in enumerate(cards):
        x = Emu(int(left + i * (cw + gap)))
        rect(s, x, top, cw, ch, fill=PANEL, rounded=True)
        rect(s, x, top, cw, Inches(0.5), fill=accent_color)  # outlet band
        text(s, Emu(int(x + Inches(0.28))), Emu(int(top + Inches(0.06))), Emu(int(cw - Inches(0.56))), Inches(0.4),
             [[(source.upper(), {"size": 11, "bold": True, "color": BG, "font": F_MONO})]], anchor=MSO_ANCHOR.MIDDLE)
        text(s, Emu(int(x + Inches(0.28))), Emu(int(top + Inches(0.75))), Emu(int(cw - Inches(0.56))), Inches(2.4),
             [[("“" + headline + "”", {"size": 18, "bold": True, "color": WHITE})]], spacing=1.1)
        text(s, Emu(int(x + Inches(0.28))), Emu(int(top + ch - Inches(0.62))), Emu(int(cw - Inches(0.56))), Inches(0.5),
             [[("▸ threatens ", {"size": 12, "color": MUTED, "font": F_MONO}),
               (tag, {"size": 12, "bold": True, "color": accent_color, "font": F_MONO})]])
    if foot: footer(s, foot)
    return s

def quote(kick, q, attrib, prompt, foot=None):
    s = slide(); kicker(s, kick)
    text(s, Inches(1.0), Inches(1.7), Inches(1.4), Inches(1.4), [[("“", {"size": 90, "bold": True, "color": GREEN})]])
    text(s, Inches(1.3), Inches(2.4), Inches(10.7), Inches(2.2), [[(q, {"size": 30, "bold": True, "color": WHITE})]], spacing=1.12)
    text(s, Inches(1.35), Inches(4.7), Inches(10.5), Inches(0.5), [[("— " + attrib, {"size": 16, "color": MUTED})]])
    if prompt:
        rect(s, Inches(1.3), Inches(5.4), Inches(10.7), Inches(1.0), fill=PANEL, line=GREEN, line_w=1.0, rounded=True)
        text(s, Inches(1.6), Inches(5.55), Inches(10.1), Inches(0.8), [[(prompt, {"size": 18, "color": WHITE})]], anchor=MSO_ANCHOR.MIDDLE)
    if foot: footer(s, foot)
    return s
