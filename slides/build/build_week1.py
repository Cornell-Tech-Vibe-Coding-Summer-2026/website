#!/usr/bin/env python3
"""Rebuild Week 1 of the 'Good Code, Good Vibes' course deck as a clean, editable .pptx.
Dark/neon course theme. All original text preserved; media shown as labeled placeholders;
thin slides lightly enriched (noted in speaker notes)."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ---- palette ---------------------------------------------------------------
BG      = RGBColor(0x0B, 0x0E, 0x14)   # near-black background
PANEL   = RGBColor(0x15, 0x1B, 0x26)   # card panels
GREEN   = RGBColor(0x00, 0xFF, 0x41)   # neon accent (brand)
WHITE   = RGBColor(0xF5, 0xF7, 0xFA)   # body text
MUTED   = RGBColor(0x8A, 0x94, 0xA6)   # captions
DIM     = RGBColor(0x5A, 0x64, 0x74)   # faint

F_HEAD = "Arial"
F_BODY = "Arial"
F_MONO = "Courier New"

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = prs.slide_width, prs.slide_height

# ---- helpers ---------------------------------------------------------------
def slide(bg=BG):
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = bg
    return s

def _set_font(run, size, color, bold, font):
    run.font.size = Pt(size); run.font.bold = bold
    run.font.color.rgb = color; run.font.name = font

def text(s, l, t, w, h, runs, size=18, color=WHITE, bold=False, font=F_BODY,
         align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, spacing=1.0, shrink=False):
    """runs: str, or list of paragraphs where each paragraph is str or list of (txt,opts)."""
    tb = s.shapes.add_textbox(l, t, w, h); tf = tb.text_frame
    tf.word_wrap = True; tf.vertical_anchor = anchor
    if isinstance(runs, str):
        runs = [runs]
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.line_spacing = spacing
        segs = para if isinstance(para, list) else [(para, {})]
        for seg_txt, opts in segs:
            r = p.add_run(); r.text = seg_txt
            _set_font(r, opts.get("size", size), opts.get("color", color),
                      opts.get("bold", bold), opts.get("font", font))
    return tb

def rect(s, l, t, w, h, fill=PANEL, line=None, line_w=1.0, rounded=False, dash=False):
    shp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE,
                             l, t, w, h)
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line; shp.line.width = Pt(line_w)
        if dash:
            shp.line._get_or_add_ln().append(_dash())
    shp.shadow.inherit = False
    return shp

def _dash():
    d = qn('a:prstDash'); from pptx.oxml import parse_xml
    el = parse_xml('<a:prstDash xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" val="dash"/>')
    return el

def kicker(s, txt, color=GREEN):
    text(s, Inches(0.7), Inches(0.55), Inches(11), Inches(0.4),
         [[(txt, {})]], size=13, color=color, bold=True, font=F_MONO)

def footer(s, txt):
    text(s, Inches(0.7), Inches(6.95), Inches(12), Inches(0.4),
         [[(txt, {})]], size=10, color=DIM, font=F_MONO)

def accent(s, l=Inches(0.7), t=Inches(1.02), w=Inches(0.9), h=Inches(0.08)):
    rect(s, l, t, w, h, fill=GREEN)

def notes(s, txt):
    s.notes_slide.notes_text_frame.text = txt

# ---- slide templates -------------------------------------------------------
def cover(title_lines, subtitle, tag):
    s = slide()
    rect(s, 0, 0, Inches(0.28), SH, fill=GREEN)
    text(s, Inches(0.9), Inches(2.2), Inches(11.5), Inches(2.2),
         [[(title_lines[0], {"size": 54, "bold": True, "color": WHITE})],
          [(title_lines[1], {"size": 54, "bold": True, "color": GREEN})]],
         spacing=1.02)
    text(s, Inches(0.92), Inches(4.5), Inches(11), Inches(0.6),
         [[(subtitle, {"size": 20, "color": MUTED})]])
    text(s, Inches(0.92), Inches(5.15), Inches(11), Inches(0.5),
         [[(tag, {"size": 13, "color": DIM, "font": F_MONO})]])
    return s

def divider(day, title, presenter=None):
    s = slide()
    rect(s, 0, 0, Inches(0.28), SH, fill=GREEN)
    text(s, Inches(0.9), Inches(2.5), Inches(11.5), Inches(0.5),
         [[(day, {"size": 16, "color": GREEN, "bold": True, "font": F_MONO})]])
    text(s, Inches(0.88), Inches(3.05), Inches(11.5), Inches(1.6),
         [[(title, {"size": 46, "bold": True, "color": WHITE})]])
    if presenter:
        text(s, Inches(0.92), Inches(4.55), Inches(11), Inches(0.5),
             [[("Led by " + presenter, {"size": 15, "color": MUTED})]])
    return s

def content(kick, title, body_paras, foot=None, title_size=34):
    s = slide(); kicker(s, kick); accent(s)
    text(s, Inches(0.7), Inches(1.3), Inches(12), Inches(1.2),
         [[(title, {"size": title_size, "bold": True, "color": WHITE})]])
    if body_paras:
        text(s, Inches(0.72), Inches(2.5), Inches(11.9), Inches(4.0),
             body_paras, size=20, color=WHITE, spacing=1.12)
    if foot: footer(s, foot)
    return s

def bullets(kick, title, items, foot=None):
    paras = [[("•  ", {"color": GREEN, "bold": True}), (it, {})] for it in items]
    return content(kick, title, paras, foot)

def big_question(kick, questions, foot=None):
    s = slide(); kicker(s, kick)
    text(s, Inches(1.2), Inches(0.6), Inches(3), Inches(2),
         [[("?", {"size": 120, "bold": True, "color": GREEN})]])
    paras = [[(q, {})] for q in questions]
    text(s, Inches(1.2), Inches(3.0), Inches(11), Inches(3.4),
         paras, size=30, color=WHITE, bold=True, spacing=1.12)
    if foot: footer(s, foot)
    return s

def media(kick, title, label, hint, foot=None):
    s = slide(); kicker(s, kick); accent(s)
    text(s, Inches(0.7), Inches(1.3), Inches(12), Inches(1.0),
         [[(title, {"size": 30, "bold": True, "color": WHITE})]])
    box = rect(s, Inches(2.4), Inches(2.55), Inches(8.5), Inches(3.9),
               fill=PANEL, line=GREEN, line_w=1.5, rounded=True, dash=True)
    text(s, Inches(2.4), Inches(3.7), Inches(8.5), Inches(0.8),
         [[(label, {"size": 22, "bold": True, "color": GREEN, "font": F_MONO})]],
         align=PP_ALIGN.CENTER)
    text(s, Inches(2.4), Inches(4.5), Inches(8.5), Inches(0.8),
         [[(hint, {"size": 14, "color": MUTED})]], align=PP_ALIGN.CENTER)
    if foot: footer(s, foot)
    return s

def columns(kick, title, cols, foot=None):
    """cols: list of (label, heading, desc)."""
    s = slide(); kicker(s, kick); accent(s)
    text(s, Inches(0.7), Inches(1.3), Inches(12), Inches(1.0),
         [[(title, {"size": 32, "bold": True, "color": WHITE})]])
    n = len(cols); gap = Inches(0.4); left = Inches(0.7)
    total = SW - Inches(1.4) - gap * (n - 1)
    cw = Emu(int(total / n))
    top = Inches(2.7); ch = Inches(3.7)
    for i, (label, head, desc) in enumerate(cols):
        x = Emu(int(left + i * (cw + gap)))
        rect(s, x, top, cw, ch, fill=PANEL, rounded=True)
        rect(s, x, top, cw, Inches(0.09), fill=GREEN)
        text(s, Emu(int(x + Inches(0.3))), Emu(int(top + Inches(0.35))),
             Emu(int(cw - Inches(0.6))), Inches(0.5),
             [[(label.upper(), {"size": 12, "bold": True, "color": GREEN, "font": F_MONO})]])
        text(s, Emu(int(x + Inches(0.3))), Emu(int(top + Inches(0.85))),
             Emu(int(cw - Inches(0.6))), Inches(0.9),
             [[(head, {"size": 20, "bold": True, "color": WHITE})]])
        text(s, Emu(int(x + Inches(0.3))), Emu(int(top + Inches(1.8))),
             Emu(int(cw - Inches(0.6))), Inches(1.7),
             [[(desc, {"size": 15, "color": MUTED})]], spacing=1.1)
    if foot: footer(s, foot)
    return s

def checklist(kick, title, items, foot=None):
    s = slide(); kicker(s, kick); accent(s)
    text(s, Inches(0.7), Inches(1.3), Inches(12), Inches(1.0),
         [[(title, {"size": 34, "bold": True, "color": WHITE})]])
    paras = [[("✓  ", {"color": GREEN, "bold": True, "size": 20}), (it, {"size": 20})]
             for it in items]
    text(s, Inches(0.9), Inches(2.5), Inches(11.5), Inches(4.2),
         paras, size=20, color=WHITE, spacing=1.35)
    if foot: footer(s, foot)
    return s

def quote(kick, q, attrib, prompt, foot=None):
    s = slide(); kicker(s, kick)
    text(s, Inches(1.0), Inches(1.7), Inches(1.4), Inches(1.4),
         [[("“", {"size": 90, "bold": True, "color": GREEN})]])
    text(s, Inches(1.3), Inches(2.4), Inches(10.7), Inches(2.2),
         [[(q, {"size": 30, "bold": True, "color": WHITE})]], spacing=1.12)
    text(s, Inches(1.35), Inches(4.7), Inches(10.5), Inches(0.5),
         [[("— " + attrib, {"size": 16, "color": MUTED})]])
    if prompt:
        rect(s, Inches(1.3), Inches(5.4), Inches(10.7), Inches(1.0),
             fill=PANEL, line=GREEN, line_w=1.0, rounded=True)
        text(s, Inches(1.6), Inches(5.55), Inches(10.1), Inches(0.8),
             [[(prompt, {"size": 18, "color": WHITE})]], anchor=MSO_ANCHOR.MIDDLE)
    if foot: footer(s, foot)
    return s

# ---- prompt-specific templates ---------------------------------------------
PROMPT_BG   = RGBColor(0x05, 0x08, 0x0C)   # near-black terminal panel
PROMPT_TXT  = RGBColor(0xD9, 0xF7, 0xE1)   # pale green-white mono text
BAD_RED     = RGBColor(0xFF, 0x5C, 0x5C)   # for "weak prompt" chips

def srcline(s, txt):
    """Small attribution line just above the footer — keeps sourcing honest."""
    text(s, Inches(0.7), Inches(6.62), Inches(12), Inches(0.35),
         [[("Source: " + txt, {})]], size=10.5, color=DIM, font=F_MONO)

def prompt_box(s, l, t, w, h, prompt_paras, label="PROMPT", chip=GREEN):
    """A visually distinct prompt: terminal-dark panel, green edge, mono text,
    labeled chip. prompt_paras: str or list of str (paragraphs)."""
    rect(s, l, t, w, h, fill=PROMPT_BG, line=chip, line_w=1.2, rounded=True)
    rect(s, l, t, Inches(0.07), h, fill=chip)                       # left edge bar
    cw = Inches(0.24 + 0.13 * len(label))
    rect(s, Emu(int(l + Inches(0.22))), Emu(int(t - Inches(0.16))), cw, Inches(0.32), fill=chip, rounded=True)
    text(s, Emu(int(l + Inches(0.22))), Emu(int(t - Inches(0.17))), cw, Inches(0.32),
         [[(label, {})]], size=11, color=PROMPT_BG, bold=True, font=F_MONO,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    if isinstance(prompt_paras, str):
        prompt_paras = [prompt_paras]
    paras = [[("> ", {"color": chip, "bold": True}), (p, {})] if i == 0 else [("  " + p, {})]
             for i, p in enumerate(prompt_paras)]
    text(s, Emu(int(l + Inches(0.32))), Emu(int(t + Inches(0.28))),
         Emu(int(w - Inches(0.6))), Emu(int(h - Inches(0.5))),
         paras, size=14, color=PROMPT_TXT, font=F_MONO, spacing=1.15)

def prompt_slide(kick, title, intro, prompt_paras, takeaway=None, src=None,
                 label="PROMPT", box_h=Inches(2.1)):
    """Title + one-liner + a single big prompt box + optional takeaway."""
    s = slide(); kicker(s, kick); accent(s)
    text(s, Inches(0.7), Inches(1.3), Inches(12), Inches(1.0),
         [[(title, {"size": 30, "bold": True, "color": WHITE})]])
    if intro:
        text(s, Inches(0.72), Inches(2.15), Inches(11.9), Inches(0.8),
             intro if isinstance(intro, list) else [[(intro, {})]],
             size=17, color=MUTED, spacing=1.1)
    prompt_box(s, Inches(0.9), Inches(3.15), Inches(11.5), box_h, prompt_paras, label=label)
    if takeaway:
        text(s, Inches(0.9), Emu(int(Inches(3.15) + box_h + Inches(0.25))), Inches(11.5), Inches(0.8),
             [[("→ ", {"color": GREEN, "bold": True}), (takeaway, {})]], size=16, color=WHITE)
    if src: srcline(s, src)
    footer(s, FOOT)
    return s

def prompt_compare(kick, title, left, right, verdict=None, src=None, box_h=Inches(2.5)):
    """Two prompts side by side. left/right: (chip_label, prompt_paras, caption)."""
    s = slide(); kicker(s, kick); accent(s)
    text(s, Inches(0.7), Inches(1.3), Inches(12), Inches(1.0),
         [[(title, {"size": 30, "bold": True, "color": WHITE})]])
    for x, (lab, paras, cap), chip in [(Inches(0.7), left, BAD_RED), (Inches(6.93), right, GREEN)]:
        prompt_box(s, x, Inches(2.55), Inches(5.7), box_h, paras, label=lab, chip=chip)
        text(s, x, Emu(int(Inches(2.55) + box_h + Inches(0.12))), Inches(5.7), Inches(0.7),
             [[(cap, {})]], size=13.5, color=MUTED, spacing=1.05)
    if verdict:
        text(s, Inches(0.72), Emu(int(Inches(2.55) + box_h + Inches(0.85))), Inches(11.9), Inches(0.7),
             [[("→ ", {"color": GREEN, "bold": True}), (verdict, {"size": 16, "color": WHITE})]])
    if src: srcline(s, src)
    footer(s, FOOT)
    return s

def term_grid(kick, title, terms, src=None, foot=None):
    """Vocabulary grid: terms = list of (term, definition), laid out 2 per row."""
    s = slide(); kicker(s, kick); accent(s)
    text(s, Inches(0.7), Inches(1.3), Inches(12), Inches(0.9),
         [[(title, {"size": 30, "bold": True, "color": WHITE})]])
    rows = (len(terms) + 1) // 2
    top = Inches(2.35); ch = Inches(1.22); gap = Inches(0.18)
    for i, (term, definition) in enumerate(terms):
        r, c = divmod(i, 2)
        x = Emu(int(Inches(0.7) + c * Inches(6.23)))
        y = Emu(int(top + r * (ch + gap)))
        rect(s, x, y, Inches(6.0), ch, fill=PANEL, rounded=True)
        rect(s, x, y, Inches(0.06), ch, fill=GREEN)
        text(s, Emu(int(x + Inches(0.25))), Emu(int(y + Inches(0.12))), Inches(5.5), Inches(0.4),
             [[(term, {})]], size=15, color=GREEN, bold=True, font=F_MONO)
        text(s, Emu(int(x + Inches(0.25))), Emu(int(y + Inches(0.52))), Inches(5.5), Emu(int(ch - Inches(0.6))),
             [[(definition, {})]], size=12.5, color=WHITE, spacing=1.05)
    if src: srcline(s, src)
    if foot: footer(s, foot)
    return s

# ---- per-day deck output ----------------------------------------------------
import os
OUT_DIR = "/Users/haukesandhaus/Documents/GitHub/Vibe-Coding-Class/website/slides"
SAVED = []

def cut_deck(filename):
    """Save the slides built so far as one per-day deck, then start a fresh one."""
    global prs, BLANK, SW, SH
    os.makedirs(OUT_DIR, exist_ok=True)
    prs.save(os.path.join(OUT_DIR, filename))
    SAVED.append((filename, len(prs.slides._sldIdLst)))
    prs = Presentation()
    prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
    BLANK = prs.slide_layouts[6]
    SW, SH = prs.slide_width, prs.slide_height

# ===========================================================================
# WEEK 1
# ===========================================================================
FOOT = "Good Code, Good Vibes · TECHIE 1121 · Cornell Tech · Summer 2026"

# --- Cover ---
cover(["Good Code, Good Vibes:", "Building Ethical Apps with AI"],
      "An introduction to ethical vibe coding", "TECHIE 1121 · CORNELL TECH · SUMMER 2026")

# ================= MONDAY — What is Vibe Coding? (Jonathan) =================
divider("WEEK 1 · MONDAY", "What is Vibe Coding?", "Jonathan")

s = quote("ETHICS WARM-UP",
      "Act only according to that maxim whereby you can at the same time will that it should become a universal law.",
      "Immanuel Kant, Groundwork of the Metaphysics of Morals",
      "If you are in poverty, is it okay to steal food for your starving family from a large grocery store?",
      FOOT)
notes(s, "Original recurring ethics-hook slide. Consider giving each day its own dilemma (kept Kant here as your opener).")

s = content("WHAT IS VIBE CODING?", "Give in to the vibes",
    [[("“There’s a new kind of coding I call ‘vibe coding’, where you fully give in to the vibes, "
       "embrace exponentials, and forget that the code even exists.”", {"size": 22, "color": WHITE})],
     [("", {})],
     [("— Andrej Karpathy, Feb 2025", {"size": 16, "color": MUTED})],
     [("", {})],
     [("You express the ", {}), ("intent", {"color": GREEN, "bold": True}),
      (" of the product; the AI does the heavy lifting. The interface is no longer "
       "syntax — it’s a conversation about behavior.", {})]], FOOT)
notes(s, "ADDED to flesh out Monday's thin 'what is vibe coding' open. Source: Karpathy / Li et al. 'Vibe Coding in Product Teams' (the 7/13 primary reading).")

s = media("WHAT IS VIBE CODING?", "See it in motion",
      "▶  vibe coding.mp4  /  whatisvibecoding.gif",
      "Drop your intro clip / GIF here (from the old deck's Drive folder).", FOOT)
notes(s, "Media placeholder — original deck had a 'what is vibe coding' gif/video here.")

content("SETTING UP", "There are many tools you can use for vibe coding",
    [[("We’ll go over how to set up and use a few of them — three different ways to work "
       "with generative AI for programming.", {"size": 22, "color": MUTED})]], FOOT)

s = columns("SETTING UP", "Three ways to work with AI",
    [("01 · Online", "Browser-only tool", "No install. Prompt and preview in the browser. e.g. v0.app."),
     ("02 · IDE", "VS Code + GitHub Copilot", "Agent mode in your editor / Codespace. (Cursor is an 18+ alternative.)"),
     ("03 · Local", "Run a model locally", "Most control over the model; add your own tools on top.")], FOOT)
notes(s, "ADDED overview to frame the three methods your Monday slides walk through. NOTE: replaced 'Antigravity' with 'VS Code + GitHub Copilot' — Antigravity is retired (no longer an IDE); Copilot is now the recommended 13+ tool. Cursor noted as 18+.")

content("METHOD 1 · ONLINE", "First, an online-only tool: v0.app",
    [[("Let’s try ", {}), ("v0.app", {"color": GREEN, "bold": True}),
      (" — you describe an app in the browser and it generates it, no setup required.", {})]], FOOT)

media("METHOD 1 · ONLINE", "v0.app — the result",
      "🖼  v0.app result — screenshot",
      "Paste your v0.app output screenshot here.", FOOT)

big_question("DISCUSS", ["Do you think this worked well?",
                         "What assumptions were made?",
                         "→ Play the game."], FOOT)

s = content("METHOD 2 · IDE", "Next, an IDE tool",
    [[("VS Code + GitHub Copilot", {"color": GREEN, "bold": True, "size": 24})],
     [("", {})],
     [("Your editor becomes the agent: it plans and edits across your whole repo. "
       "Cursor is a similar option (18+).", {"size": 20, "color": WHITE})]], FOOT)
notes(s, "Original listed 'Antigravity / VSCode / Cursor'. Antigravity retired → now VS Code + GitHub Copilot (recommended, 13+); Cursor kept as 18+ alt.")

content("METHOD 2 · IDE", "Let’s walk through VS Code first",
    [[("Install VS Code, add the GitHub Copilot extension, open your repo, and turn on "
       "agent mode.", {"size": 22, "color": MUTED})]], FOOT)

media("METHOD 2 · IDE", "Once it’s installed it looks like this",
      "▶  vscode.mp4",
      "Drop your VS Code walkthrough clip / screenshot here.", FOOT)

big_question("DISCUSS", ["What are the benefits of the code living right next to the generative AI?",
                         "What are the negatives?"], FOOT)

content("METHOD 3 · LOCAL", "Lastly, run a model locally",
    [[("This gives you the ", {}), ("most control", {"color": GREEN, "bold": True}),
      (" over the AI models themselves — and you can add many different tools on top of the model.",
       {})]], FOOT)

media("METHOD 3 · LOCAL", "Running a model locally",
      "▶  local.mp4",
      "Drop your local-model clip / screenshot here.", FOOT)

big_question("DISCUSS", ["What could you do running the model locally that you can’t with the other two methods?",
                         "If running the program is separated from the model, what problems can that cause?"], FOOT)

big_question("DISCUSS", ["Given the three methods of interacting with generative AI for programming —",
                         "which method is best in which situation?"], FOOT)

checklist("GOALS FOR THE DAY", "Goals for the day",
    ["GitHub repo cloned locally to your computer",
     "Updates pushed to your own repo copy",
     "Set up and chosen a primary vibe coding tool",
     "Worked on Activity 1 (your website)",
     "Website is hosted",
     "Submitted your Vibe-Trace",
     "Submitted your activity reflection"], FOOT)

# ================= TUESDAY — Prompt Engineering (Hauke) =================
cut_deck("Week1-Mon-What-Is-Vibe-Coding.pptx")
divider("WEEK 1 · TUESDAY", "Prompt Engineering", "Hauke")

# --- Framing: the day's axis ---
s = content("TODAY", "One axis: vibe-based vs. engineered prompting",
    [[("Vibe coding leans on intuition and natural language. ", {"size": 22, "color": WHITE}),
      ("Prompt engineering", {"size": 22, "color": GREEN, "bold": True}),
      (" is the disciplined version — closer to writing software requirements.", {"size": 22, "color": WHITE})],
     [("", {})],
     [("You’ll feel the difference twice today: fast (the four challenges), then for real "
       "(a deliberate 25+25-minute fork of your own project).", {"size": 20, "color": MUTED})]], FOOT)
notes(s, "Framing mirrors the 7/14 instructions ('Prompts that Steer the Vibe'). Everything on today's slides carries a source line — we practice the attribution we preach.")

# --- The reading: Why Johnny Can't Prompt ---
s = content("READING", "Why Johnny Can’t Prompt (CHI ’23)",
    [[("The study: 10 people with little or no prompting experience, given a no-code tool "
       "for building a chatbot purely with prompts.", {"size": 22, "color": WHITE})],
     [("", {})],
     [("Everyone could prompt. Almost nobody could prompt ", {"size": 22, "color": WHITE}),
      ("systematically", {"size": 22, "color": GREEN, "bold": True}),
      (" — and the ways they failed are exactly the ways you’ll fail this afternoon.", {"size": 22, "color": WHITE})],
     [("", {})],
     [("Zamfirescu-Pereira, Wong, Hartmann & Yang — today’s primary reading.", {"size": 15, "color": MUTED})]], FOOT)
notes(s, "Design probe: BotDesigner (GPT-3 chatbot design tool). Participants across ages/professions, incl. a software engineer — the failure modes were NOT about programming skill.")

s = bullets("READING · FINDINGS", "How non-experts actually prompt",
    ["Opportunistic tinkering — change something, retry, repeat. No hypothesis, no record.",
     "Over-generalizing from one try: one failure → “it can’t do that”; one success → declare victory.",
     "Human-social habits: staying polite, avoiding repetition — and refusing to paste examples because it “feels like cheating” (it isn’t; it works).",
     "Reaching for “don’t do X” — even though “do Y” works far better.",
     "Expecting human capabilities: memory across chats, knowing who is being addressed.",
     "And eventually, every participant asked: “but why did it do that?”"], FOOT)
srcline(s, "Zamfirescu-Pereira, Wong, Hartmann & Yang, “Why Johnny Can’t Prompt”, CHI ’23 — dl.acm.org/10.1145/3544548.3581388")
notes(s, "Each bullet is a finding from §4 of the paper. The 'cheating' quote and do/don't asymmetry are verbatim study observations (P2, P4, P9).")

s = prompt_compare("READING · IN PRACTICE", "Say what you DO want",
    ("DON’T-STYLE", "Do not use markdown in your response.",
     "Negations often just don’t land. In the study, P4 stacked “don’t say…” instructions that had no effect at all."),
    ("DO-STYLE", "Your response should be composed of smoothly flowing prose paragraphs.",
     "Positive framing states the target behavior — nothing to mis-negate."),
    verdict="Tell it what to do, not what to avoid — and explain why, so the model can generalize.",
    src="Anthropic, Claude prompting best practices · Zamfirescu-Pereira et al., CHI ’23", box_h=Inches(1.9))
notes(s, "The paper: interviewers observed 'Do not do X' much less effective than 'Do Y', yet participants overwhelmingly chose the former. Early-childhood-education literature shows the same bias in humans.")

s = checklist("THE DISCIPLINE", "Prompt like an engineer, not a tinkerer",
    ["Know why the last prompt failed before writing the next one.",
     "Change ONE variable at a time — wording, example, or structure.",
     "Test more than once before concluding “it can’t do that”.",
     "Use examples. It isn’t cheating — it’s the most reliable steering tool there is.",
     "Keep the receipts: your auto-logged history.md is your lab notebook."], FOOT)
srcline(s, "iteration discipline from Zamfirescu-Pereira et al., CHI ’23 (§5.1)")
notes(s, "This checklist IS the Activity 2 run-2 rubric. The paper's core training implication: 'end-users should collect more data than they are naturally inclined to.'")

# --- Anatomy & vocabulary ---
s = prompt_slide("ANATOMY", "The four parts of a prompt",
    "Most prompts contain up to four elements. Label them mentally and you can debug them separately.",
    ["Instruction:  Build a one-page portfolio site in HTML/CSS.",
     "Context:      I’m a photography student; monochrome, image-first.",
     "Input data:   Use the bio and six image captions pasted below.",
     "Output:       A single index.html, no frameworks, mobile-first."],
    takeaway="When a prompt fails, ask which element was missing or vague — instruction, context, input, or output format.",
    src="element framework: DAIR.AI Prompt Engineering Guide — promptingguide.ai/introduction/elements (MIT)",
    label="ANATOMY", box_h=Inches(2.2))
notes(s, "DAIR.AI's canonical framework (instruction/context/input data/output indicator), recast with a course-shaped example.")

s = term_grid("VOCABULARY · 1/2", "Words you need this week",
    [("prompt", "Everything the model sees this turn — instruction, context, examples, data."),
     ("system prompt", "Standing instructions that frame the whole session and outrank the chat — where roles and rules live."),
     ("token", "The unit models read and write (≈¾ of a word). You pay, and wait, per token."),
     ("context window", "The model’s working memory. Files, chat, and output must all fit — overflow gets forgotten."),
     ("temperature", "The randomness dial: low = precise and repeatable, high = creative and loose."),
     ("hallucination", "Fluent, confident output that is simply wrong. Grounding in real files and docs is the antidote.")],
    src="definitions after promptingguide.ai & Anthropic docs", foot=FOOT)

s = term_grid("VOCABULARY · 2/2", "Words you need this week",
    [("zero-shot", "Just ask, no examples — the default mode with modern instruction-tuned models."),
     ("few-shot", "Show 2–5 examples of what you want first; the model copies the pattern."),
     ("chain-of-thought", "Make the model reason before answering — “explain your plan before writing code”."),
     ("prompt chaining", "Split a big task into steps; each prompt builds on the last output. Debuggable."),
     ("RAG / grounding", "Hand the model real documents to answer from instead of memory — @-files, pasted docs."),
     ("agent (ReAct loop)", "Think → act (run a tool) → observe → repeat. This is what Copilot agent mode does.")],
    src="definitions after promptingguide.ai & Anthropic docs", foot=FOOT)

# --- Principles (UVic) ---
s = bullets("PRINCIPLES", "Seven principles of prompt engineering",
    ["Be concise — simple language, no filler.",
     "Be clear — no vague or ambiguous wording.",
     "Include context & logical structure.",
     "Break down complex tasks — one ask per prompt.",
     "Specify the output — style, tone, depth, length, format.",
     "Reflect & adapt — not the outcome you wanted? Adjust, re-run.",
     "Combine them — the principles compound."], FOOT)
srcline(s, "UVic Libraries, “Prompt Engineering for GenAI — Beginner Course” (CC BY 4.0) — libguides.uvic.ca/Prompt_Engineering_Beginners_Course")
notes(s, "UVic's 7 principles, lightly compressed. CC BY 4.0 — attribution on-slide satisfies the license.")

s = prompt_compare("PRINCIPLE · SPECIFY OUTPUT", "Every unstated detail is a decision you delegate",
    ("VAGUE", "Draft a paragraph on climate change.",
     "No audience, style, length, or format — the model picks its averaged defaults for all four."),
    ("SPECIFIED", "Draft a newspaper paragraph of 500 words including a headline, on the effects of climate change. Apply the simplified grammar of a tabloid paper, and use sensational language.",
     "Style, tone, length, and format all pinned down — one shot lands far closer."),
    verdict="Challenge C is exactly this experiment: 30 seconds of vibes vs. 10 minutes of spec, one shot each.",
    src="example prompts: UVic Libraries (CC BY 4.0)", box_h=Inches(2.5))

s = prompt_compare("PRINCIPLE · BE DIRECT", "If you want ambition, ask for it",
    ("BARE", "Create an analytics dashboard",
     "Correct but minimal — the model does exactly this, and nothing more."),
    ("EXPLICIT", "Create an analytics dashboard. Include as many relevant features and interactions as possible. Go beyond the basics to create a fully-featured implementation.",
     "“Above and beyond” only happens when you request it."),
    verdict="Golden rule: if a colleague with no context would be confused by your prompt, the model will be too.",
    src="Anthropic, Claude prompting best practices — platform.claude.com/docs", box_h=Inches(2.5))

# --- Techniques ---
s = prompt_slide("TECHNIQUE · ROLE", "Give the model a role",
    "A persona changes code quality, not just tone: naming, comments, error handling. You’ll test this in Challenge B.",
    ["You are a senior software engineer. Write a clean, PEP8-compliant",
     "grade calculator with docstrings, error handling for non-numeric",
     "input, and descriptive variable names."],
    takeaway="Even one sentence of role shifts every default. UVic’s framing: “You are the movie director — the AI is your method actor.”",
    src="UVic Libraries (CC BY 4.0) · Anthropic best practices (role prompting)", label="ROLE PROMPT", box_h=Inches(1.9))

s = prompt_slide("TECHNIQUE · EXAMPLES", "Show, don’t tell (few-shot)",
    [[("The most reliable way to steer format, tone, and structure — and the one the Johnny study’s "
       "participants avoided because it “felt like cheating.” Use 3–5 relevant, diverse examples.", {})]],
    ["<examples>",
     "  <good> …a card component in exactly the style you want… </good>",
     "  <bad> …the over-decorated version to avoid… </bad>",
     "</examples>",
     "Write the remaining three cards in the style of the good example."],
    takeaway="In Activity 2, one good + one bad snippet counts as a toolkit technique.",
    src="Anthropic best practices (multishot) · promptingguide.ai/techniques/fewshot", label="FEW-SHOT", box_h=Inches(2.3))

s = prompt_slide("TECHNIQUE · STRUCTURE", "Mark up your prompt (XML tags)",
    "When role, style, and requirements blur into one paragraph, the model guesses at the boundaries. Tags make them explicit — and turn your prompt into a reusable template.",
    ["<role>Expert frontend developer</role>",
     "<style_guide>Dark mode, high contrast, neon borders</style_guide>",
     "<requirements>",
     "- Sidebar with 4 nav links",
     "- Central stats grid",
     "</requirements>"],
    takeaway="This is the exact pattern Activity 2 asks for in your spec-first prompt.",
    src="Anthropic best practices (XML tags) · UVic technique 5, “mark-up your input” (CC BY 4.0)", label="XML TAGS", box_h=Inches(2.5))

s = prompt_slide("TECHNIQUE · THINK FIRST", "Make it plan before it codes",
    "Modern agents reason by default — your job is to steer it and make the plan visible so you can check it before code exists. Prefer general nudges (“think thoroughly”) over prescriptive step lists.",
    ["Before writing any code, explain your layout logic inside",
     "<thinking> tags.",
     "Before you finish, verify the result against the requirements above."],
    takeaway="A wrong plan is cheaper to fix than wrong code. This is Activity 2’s “thinking space” technique.",
    src="Anthropic best practices (extended thinking, self-check) · promptingguide.ai/techniques/cot", label="PLAN FIRST", box_h=Inches(1.9))

s = prompt_slide("TECHNIQUE · CHAINING", "One ask per prompt — then chain them",
    "Overloaded prompts fail silently: five asks, three delivered, and you don’t notice. Chains are debuggable — you see exactly which step went wrong.",
    ["1. Scaffold the page structure — HTML only, no styling yet.",
     "2. Style it, following the <style_guide> from before.",
     "3. Add the API call, with loading and error states.",
     "4. Now write tests; fix whatever fails."],
    takeaway="This is how experienced vibe coders actually work — the star technique for coding agents.",
    src="promptingguide.ai/techniques/prompt_chaining · UVic principle 4 (CC BY 4.0)", label="PROMPT CHAIN", box_h=Inches(2.3))

s = columns("UNDER THE HOOD", "What your agent does between your prompts",
    [("ReAct loop", "Think → act → observe",
      "It reasons, runs a tool (read a file, run tests), observes the result, reasons again. Letting it run tests makes it smarter."),
     ("RAG / grounding", "Retrieve, then answer",
      "Answers improve when the model pulls in real context — your files, docs, README — instead of guessing from memory. That’s why @-mentioning files works."),
     ("Self-correction", "Draft → review → refine",
      "The most useful chain of all: generate, have it critique its own output against your criteria, then revise.")], FOOT)
srcline(s, "promptingguide.ai (ReAct, RAG) · Anthropic best practices (self-correction chains)")
notes(s, "Students don't write ReAct prompts — they should recognize the loop to understand what agent mode is doing and why tool access matters.")

# --- Frontend aesthetics (maps to Challenges A & D) ---
s = prompt_slide("TECHNIQUE · AESTHETICS", "Escaping the “AI slop” look",
    "Models converge on statistically safe design: Inter, purple gradients, cookie-cutter layouts. Anthropic’s own frontend guidance tells the model to fight that — steal the trick.",
    ["Avoid generic AI aesthetics: overused fonts (Inter, Roboto, Arial),",
     "cliched purple-gradient-on-white schemes, predictable layouts.",
     "Commit to a cohesive theme: distinctive typography, dominant colors",
     "with sharp accents, layered backgrounds, one well-orchestrated",
     "page-load animation. Make choices designed for THIS context."],
    takeaway="Name the aesthetic and the audience: “a 1980s neon arcade at midnight” beats “20px border radius” — that’s Challenge D.",
    src="condensed from Anthropic’s frontend-design guidance — platform.claude.com/docs (#frontend-design)", label="AESTHETICS", box_h=Inches(2.3))
notes(s, "Condensed from Anthropic's <frontend_aesthetics> system-prompt snippet (docs, #frontend-design) and the frontend-design skill. Full snippet is worth sharing as a handout.")

s = media("TECHNIQUE · AESTHETICS", "Same prompt, with vs. without design guidance",
      "🖼  before / after pairs — Anthropic blog",
      "Drop in before/after screenshots from “Improving frontend design through Skills” (URLs in speaker notes).", FOOT)
notes(s, "Blog: https://claude.com/blog/improving-frontend-design-through-skills — best pairs:\n"
         "Blog layout without/with skill: ...6913d5b728dcecc13bc1f78d_f7040147.png / ...6913d5b728dcecc13bc1f77e_0ce357ff.png\n"
         "Dashboard without/with: ...6913d5b728dcecc13bc1f784_7beb17d0.png / ...6913d5b728dcecc13bc1f781_3705adad.png\n"
         "Task manager without/with: ...6913d5b728dcecc13bc1f793_875d1eef.png / ...6913d5b728dcecc13bc1f7c9_7ae52606.png\n"
         "All under https://cdn.prod.website-files.com/68a44d4040f98a4adf2207b6/ — credit the blog on-slide.")

s = bullets("TECHNIQUE · AESTHETICS", "impeccable.style — design vocabulary, packaged",
    ["An opinionated, open-source skill you install next to your agent — “the missing design vocabulary for agents” (Paul Bakaus).",
     "23 design commands (/typeset, /colorize, /animate, …) so you can ask for the exact thing without explaining it.",
     "46 deterministic “anti-slop” rules that strip the AI-default design tells automatically.",
     "Works with Copilot, Claude Code, Cursor, and more — use it as your reference when a design feels generic.",
     "impeccable.style"], FOOT)
srcline(s, "impeccable.style — Paul Bakaus (free & open source)")
notes(s, "Reference library for students: instead of hand-writing anti-slop prompts every time, install the skill (or crib its vocabulary). Pairs with the previous two slides.")

s = prompt_compare("AGENT PROMPTS", "Say “change it”, not “could you suggest…”",
    ("SUGGESTS", "Can you suggest some changes to improve this function?",
     "The agent describes improvements — and touches nothing."),
    ("ACTS", "Change this function to improve its performance.",
     "The agent edits the code. Your verbs decide: advice or action."),
    verdict="More agent guardrails that pay off: “read the file before answering” · “minimum complexity needed” · “clean up temp files when done.”",
    src="Anthropic, Claude prompting best practices (coding agents)", box_h=Inches(1.9))

# --- Hauke's field notes (kept from the original tips) ---
s = content("FIELD NOTES · 01", "Make it personal",
    [[("Generic prompt → generic app — the model falls back on its averaged defaults "
       "(remember the Mirror Effect).", {"size": 20, "color": WHITE})],
     [("", {})],
     [("Feed it ", {}), ("your", {"color": GREEN, "bold": True}),
      (" specifics: your content, your voice, real names and photos, a reference you "
       "actually like. Concrete context is what pulls the output away from the default.", {})]], FOOT)
notes(s, "Hauke's own tips, kept from the original deck ('Nail the look' folded into the aesthetics slide).")

content("FIELD NOTES · 02", "Vibe-code the controls, not just the result",
    [[("Don’t re-prompt for every tweak — ask the AI to expose knobs:", {"size": 20, "color": WHITE})],
     [("", {})],
     [("“Add sliders / CSS variables for spacing, colors, and font size so I can fine-tune it myself.”",
       {"size": 20, "color": GREEN, "bold": True})],
     [("", {})],
     [("Now you iterate by hand — fast, and without spending another generation.", {"size": 20, "color": MUTED})]], FOOT)

bullets("FIELD NOTES · 03", "Deal with quota",
    ["Free tiers run out — every wasted regeneration is wasted quota.",
     "Plan the prompt before you spend a generation.",
     "Make one bigger, specific ask instead of ten vague ones.",
     "Save your transcripts, and know when to switch tools."], FOOT)

content("FIELD NOTES · 04", "The value of a vibe-coded artifact",
    [[("What makes it good isn’t that AI made it — it’s ", {"size": 20, "color": WHITE}),
      ("your taste and iteration", {"size": 20, "color": GREEN, "bold": True}),
      (".", {"size": 20, "color": WHITE})],
     [("", {})],
     [("e.g. Isometric NYC — cannoneyed.com/projects/isometric-nyc — a striking artifact that "
       "reflects human judgment, not a one-click default.", {"color": MUTED})],
     [("", {})],
     [("Aim for something only you would have made.", {})]], FOOT)

# --- Activities ---
s = columns("ACTIVITY 1", "Four challenges — pick one, feel the axis",
    [("A · Timer", "Mood vs. blueprint",
      "A Pomodoro timer: “rainy coffee shop” vs. explicit requirements. → specify the output"),
     ("B · Persona", "Role prompting",
      "A grade calculator, ± “you are a senior engineer.” → role prompting"),
     ("C · Spec density", "30 sec vs. 10 min",
      "A link-in-bio page: throwaway prompt vs. dense spec. → be clear & specific"),
     ("D · Vaporwave", "Aesthetic language",
      "A game login screen: CSS specs vs. “neon arcade at midnight.” → aesthetics")], FOOT)
notes(s, "Run both prompts once each, in FRESH chats — one shot, no fixing. We share right after: one sentence on where vibes won and where engineering was non-negotiable.")

s = bullets("ACTIVITY 2", "Fork your brain (25 + 25 minutes)",
    ["Build the same project twice — once vibes-only, once engineered. The timer IS the experiment.",
     "Run 1 — engineering brain OFF: just talk, chase the vibe, stop when the timer rings.",
     "Run 2 — spec-first prompt using ≥2 toolkit techniques: XML tags, thinking space, or few-shot.",
     "Iterate like Johnny couldn’t: know why the last prompt failed, change one variable at a time.",
     "No note-taking needed — history.md logs every prompt; you’ll analyze it in the vibe report."], FOOT)
notes(s, "Pick something genuinely yours (a tool you'd use, an obsession page, a gift). Equal time is what makes the versions comparable — no 'five more minutes'.")

# --- Sources & going deeper ---
s = content("SOURCES", "Credits — and where to go deeper",
    [[("UVic Libraries", {"color": GREEN, "bold": True}),
      (" — Prompt Engineering for GenAI, beginner course (CC BY 4.0): the seven principles & example prompts. libguides.uvic.ca", {"size": 17})],
     [("DAIR.AI Prompt Engineering Guide", {"color": GREEN, "bold": True}),
      (" — prompt anatomy, techniques, tips (MIT). promptingguide.ai — the deep well when you want more.", {"size": 17})],
     [("Anthropic", {"color": GREEN, "bold": True}),
      (" — Claude prompting best practices & frontend-design guidance (platform.claude.com/docs); hands-on 9-chapter tutorial: github.com/anthropics/prompt-eng-interactive-tutorial", {"size": 17})],
     [("Zamfirescu-Pereira, Wong, Hartmann & Yang", {"color": GREEN, "bold": True}),
      (" — “Why Johnny Can’t Prompt”, CHI ’23: the study behind today’s discipline.", {"size": 17})],
     [("impeccable.style", {"color": GREEN, "bold": True}),
      (" — Paul Bakaus’s open-source design vocabulary & anti-slop rules for coding agents.", {"size": 17})]], FOOT)
notes(s, "Per-slide source lines + this credits slide keep the borrowing honest. The Anthropic tutorial is the recommended self-study path (note: its 'prefill' chapter is deprecated on current Claude models).")

# ================= WEDNESDAY — Bias in Vibe Coding (Jonathan) =============
cut_deck("Week1-Tue-Prompt-Engineering.pptx")
divider("WEEK 1 · WEDNESDAY", "Bias in Vibe Coding", "Jonathan")

columns("LEARNING OBJECTIVES", "Learning objectives",
    [("Visibility", "Visual & functional bias",
      "Identify how LLM training data manifests as visual and functional bias in generated apps."),
     ("Analytics", "Tool evaluation",
      "Evaluate different vibe coding tools (Lovable, Bolt, Vibe Code App) for consistency and stereotyping."),
     ("Auto-awesome", "Inclusive design",
      "Develop strategies for ‘de-biasing’ prompts to ensure more inclusive outputs.")], FOOT)

content("FOUNDATION", "What is bias?",
    [[("Bias is a systematic skew in outputs — patterns the model treats as the default "
       "because they dominate its training data.", {"size": 22, "color": MUTED})]], FOOT)
notes(prs.slides[-1], "Original 'What is Bias?' slide was a title only — added a one-line working definition.")

columns("DEFINITION", "What is “Vibe Bias”?",
    [("Visual stereotypes", "Palette & aesthetics",
      "e.g. ‘Professional’ defaults to blue, white, corporate aesthetics."),
     ("Functional assumptions", "Defaults baked in",
      "e.g. assuming US-style addresses or phone-number formats."),
     ("Accessibility oversights", "Who’s left out",
      "e.g. low contrast ratios or non-semantic HTML structures.")], FOOT)
notes(prs.slides[-1], "Kept your 'Vibe Bias' definition; the lead definition line: 'The tendency of LLMs to generate UI/UX patterns based on the most frequent (often Western-centric) data in their training sets.'")

content("CORE CONCEPT", "The Mirror Effect",
    [[("LLMs don’t “design” — they predict the next most likely token.", {"size": 24, "color": WHITE, "bold": True})],
     [("", {})],
     [("Prompt for a “Healthcare Dashboard” and the tool reflects the average of 10,000 "
       "existing dashboards — including all their flaws and lack of diversity.", {"size": 20, "color": MUTED})]], FOOT)

big_question("TRANSITION", ["Time for some activities."], FOOT)

s = slide(); kicker(s, "CASE STUDY"); accent(s)
text(s, Inches(0.7), Inches(1.3), Inches(12), Inches(1.0),
     [[("Case study: “Make it Professional”", {"size": 32, "bold": True, "color": WHITE})]])
rect(s, Inches(0.7), Inches(2.6), Inches(5.7), Inches(2.6), fill=PANEL, rounded=True, line=GREEN, line_w=1.2, dash=True)
text(s, Inches(0.7), Inches(3.5), Inches(5.7), Inches(0.6),
     [[("🖼  Result A", {"size": 18, "bold": True, "color": GREEN, "font": F_MONO})]], align=PP_ALIGN.CENTER)
text(s, Inches(0.7), Inches(4.1), Inches(5.7), Inches(0.6),
     [[("“A standard landing page.”", {"size": 14, "color": MUTED})]], align=PP_ALIGN.CENTER)
rect(s, Inches(6.9), Inches(2.6), Inches(5.7), Inches(2.6), fill=PANEL, rounded=True, line=GREEN, line_w=1.2, dash=True)
text(s, Inches(6.9), Inches(3.5), Inches(5.7), Inches(0.6),
     [[("🖼  Result B", {"size": 18, "bold": True, "color": GREEN, "font": F_MONO})]], align=PP_ALIGN.CENTER)
text(s, Inches(6.9), Inches(4.1), Inches(5.7), Inches(0.6),
     [[("“A professional landing page.”", {"size": 14, "color": MUTED})]], align=PP_ALIGN.CENTER)
text(s, Inches(0.7), Inches(5.5), Inches(12), Inches(0.8),
     [[("The question: why does “professional” often default to one specific layout or design?",
        {"size": 18, "color": WHITE, "bold": True})]])
footer(s, FOOT)

bullets("ACTIVITY", "The cross-tool benchmark",
    ["The mission: run the SAME prompt across three different platforms.",
     "The prompt: “Build a social networking profile page for a community organizer in a "
     "specific area you’re familiar with.”",
     "Look for: did it use relevant imagery or icons?",
     "Is the terminology correct for the location?",
     "Does the wording or layout assume anything about who uses the site?"], FOOT)

columns("AUDIT", "The Vibe Audit checklist",
    [("Representation", "Who is this for?",
      "Who is this interface designed for? Who might be left out or overlooked?"),
     ("Accessibility", "Does it meet the bar?",
      "Color contrast, readability, clear navigation — basic accessibility standards."),
     ("Assumptions", "What’s baked in?",
      "What does the design assume — payment access, name formats, language?")], FOOT)

columns("STRATEGIES", "Strategies for control",
    [("Explicit prompting", "Say what you mean",
      "From “Make a login” to “Make an accessible login using semantic HTML and high-contrast tokens.”"),
     ("Context loading", "Feed diverse refs",
      "Provide the AI with diverse design-system references before generating."),
     ("Human-in-the-loop", "Audit, don’t trust",
      "Use your own critical lens to audit the code — not just the ‘vibe.’")], FOOT)

# ================= THURSDAY — Project Day =================
cut_deck("Week1-Wed-Bias-In-Vibe-Coding.pptx")
divider("WEEK 1 · THURSDAY", "Project Day — Project 1")

s = bullets("PROJECT 1", "Vibe code something your team wants or needs",
    ["Teams of 3–4. Everyone codes — rotate who drives the AI.",
     "Pick a real problem your team has and ship a small hosted app for it.",
     "Use the strongest tool + prompt approach you found Mon–Wed.",
     "Deliverable: hosted link + a vibe report (tool choice, prompt log, what broke).",
     "Work happens in your team’s Group Project repo. Due Monday of Week 2, before class."], FOOT)
notes(s, "Enriched the empty 'Project Day' slide with the Project 1 brief essentials from the repo. Group projects are separate Classroom team repos (confirm details with Jonathan).")

# ---------------------------------------------------------------------------
cut_deck("Week1-Thu-Project-Day.pptx")
for name, n in SAVED:
    print(f"Saved {name} · {n} slides")
