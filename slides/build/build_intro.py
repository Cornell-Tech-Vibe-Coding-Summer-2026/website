#!/usr/bin/env python3
"""Course intro / Week 0 — cover, agenda, rhythm, setup, instructors (with links to their
vibe-coded pages), what you'll learn, the values we embed, how the course was itself
vibe-coded (colophon), and an honest 'limits of our ethics' slide."""

from deck_common import *
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
prs = init_deck()

FOOT = "Good Code, Good Vibes · TECHIE 1121 · Cornell Tech · Summer 2026"

HAUKE_PAGE = "https://vibe-coding-ethics.tech.cornell.edu/week1/7_13/examples/activity1_basic-example/code_deliverable/"
JON_PAGE   = "https://vibe-coding-ethics.tech.cornell.edu/week1/7_13/examples/activity1-jonathan/code_deliverable/"
COLOPHON   = "https://cornell-tech-vibe-coding-summer-2026.github.io/website/?lite#notepad"

def linked(s, l, t, w, label, url, size=15, color=GREEN):
    tb = s.shapes.add_textbox(l, t, w, Inches(0.4)); tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; r = p.add_run(); r.text = label
    r.font.size = Pt(size); r.font.bold = True; r.font.color.rgb = color; r.font.name = F_MONO
    r.hyperlink.address = url
    return tb

# 1 — Cover
cover(["Good Code, Good Vibes:", "Building Ethical Apps with AI"],
      "An introduction to ethical vibe coding", "TECHIE 1121 · CORNELL TECH · SUMMER 2026")

# 2 — Agenda
columns("AGENDA", "Three weeks",
    [("Week 1 · Gaining Control", "Vibe Coding",
      "What is vibe coding · prompt engineering · bias in vibe coding · Project 1."),
     ("Week 2 · Doing the Right Thing", "Ethics",
      "Values · AI safety & red-teaming · AI-against-AI · Project 2."),
     ("Week 3 · Useful & Empowering", "Outcomes",
      "Usability testing · values testing · final project & presentations.")], FOOT)

# 3 — Rhythm & deliverables
bullets("HOW IT WORKS", "The rhythm — and what you hand in",
    ["Mon–Thu for three weeks (Jul 13–30).",
     "Each week: three hands-on activity days + a Thursday group-project day.",
     "Keep a Vibe-Trace — your AI prompts are logged as you build. It’s graded.",
     "Repo-only: your GitHub repo is your submission, and it deploys to a live page.",
     "Three group projects build toward a final that helps someone outside the class."], FOOT)

# 4 — Setup
s = bullets("GET SET UP", "Before we start",
    ["Sign up for GitHub, then join the Student Developer Pack — free GitHub Copilot (13+).",
     "Accept the course GitHub Classroom assignment — that becomes your repo.",
     "Open it in VS Code (Copilot extension) or a Codespace, and turn on agent mode.",
     "Under 18? Use Copilot. Claude Code and Cursor are 18+ — the repo works with all of them."], FOOT)
notes(s, "Setup mirrors the current 7_13 activity. Jonathan may have tweaked the live version — reconcile with his slide if so.")

# 5 — Instructors
s = slide(); kicker(s, "MEET YOUR INSTRUCTORS"); accent(s)
text(s, Inches(0.7), Inches(1.3), Inches(12), Inches(1.0),
     [[("Who’s teaching this — and eating our own cooking", {"size": 30, "bold": True, "color": WHITE})]])
for x, name, role, page in [
    (Inches(0.7), "Hauke Sandhaus", "Ethics & design researcher, Cornell Tech.", HAUKE_PAGE),
    (Inches(6.93), "Jonathan Segal", "PhD researcher, Cornell Tech — AR/XR for healthcare & EMS.", JON_PAGE)]:
    rect(s, x, Inches(2.7), Inches(5.7), Inches(3.2), fill=PANEL, rounded=True)
    rect(s, x, Inches(2.7), Inches(5.7), Inches(0.09), fill=GREEN)
    text(s, Inches(x.inches + 0.35), Inches(3.05), Inches(5.0), Inches(0.6),
         [[(name, {"size": 24, "bold": True, "color": WHITE})]])
    text(s, Inches(x.inches + 0.35), Inches(3.75), Inches(5.0), Inches(1.0),
         [[(role, {"size": 16, "color": MUTED})]], spacing=1.1)
    linked(s, Inches(x.inches + 0.35), Inches(5.15), Inches(5.0),
           "▶  their vibe-coded “about me” page", page, size=15)
text(s, Inches(0.72), Inches(6.2), Inches(12), Inches(0.4),
     [[("With Wendy Ju (faculty). These pages are our own Week-1 activity — your inspiration to steal from.",
        {"size": 13, "color": DIM})]])
footer(s, FOOT)
notes(s, "Links go to Hauke's and Jonathan's hosted, vibe-coded 'About Me' pages (the 7/13 example). Real sites: haukesandhaus.de · jonathansegal.io.")

# 6 — What you'll learn
content("LEARNING OUTCOME", "What you’ll learn",
    [[("Use AI coding tools to build real artifacts — ", {"size": 24, "color": WHITE}),
      ("ethically", {"size": 24, "color": GREEN, "bold": True}),
      (" and ", {"size": 24, "color": WHITE}),
      ("effectively", {"size": 24, "color": GREEN, "bold": True}), (".", {"size": 24, "color": WHITE})],
     [("", {})],
     [("Ethically", {"color": GREEN, "bold": True}),
      (" — conscientiously: name the values, test with real people, own your calls.", {})],
     [("Effectively", {"color": GREEN, "bold": True}),
      (" — as a good collaborator: steer the AI, review its work, and stay in charge of the outcome.", {})]], FOOT)

# 7 — Values we embed
content("OUR COMMITMENTS", "The values we embed in this course",
    [[("Transparency", {"color": GREEN, "bold": True}),
      (" — we show our work, and you keep a Vibe-Trace.", {})],
     [("Honing human capability", {"color": GREEN, "bold": True}),
      (" — AI to sharpen your judgment, not replace it.", {})],
     [("Non-manipulation", {"color": GREEN, "bold": True}),
      (" — no dark patterns; honest, humane design.", {})],
     [("A trustworthy web", {"color": GREEN, "bold": True}),
      (" — real sources, accessible, no AI slop passed off as fact.", {})]], FOOT)

# 8 — We vibe-coded this course too
s = content("PRACTICE WHAT WE TEACH", "We vibe-coded this course, too",
    [[("This course and its 3D site were themselves vibe-coded over ~6 months — with Claude, "
       "GitHub Copilot, Gemini, and Warp; the 3D figures via Meshy AI; the scene in Blender.",
       {"size": 21, "color": WHITE})],
     [("", {})],
     [("Every AI turn is logged with the same Vibe-Trace we ask of you, and both repos are public. "
       "We hold ourselves to the standard we set for you.", {"size": 19, "color": MUTED})]], FOOT)
linked(s, Inches(0.72), Inches(5.4), Inches(11),
       "▶  read the full AI-Use Disclosure & Colophon", COLOPHON, size=15)
notes(s, "Colophon link deep-links to the disclosure on the showcase site (lite notepad).")

# 9 — Limits of our ethics
content("BE HONEST", "The limits of our ethics",
    [[("We’re upfront about the frame:", {"size": 22, "color": WHITE})],
     [("", {})],
     [("•  ", {"color": GREEN, "bold": True}),
      ("We take vibe coding as a practice — we teach how to do it conscientiously, not whether "
       "it should exist.", {})],
     [("•  ", {"color": GREEN, "bold": True}),
      ("Our lens assumes a human designer is still in charge. Critiques of that assumption — "
       "power, labor, the environment, who gets to build — sit largely outside this short course.", {})],
     [("•  ", {"color": GREEN, "bold": True}),
      ("“Conscientious design” is a starting point, not the last word on AI ethics.", {})]], FOOT)
notes(prs.slides[-1], "Hauke's framing: learning outcome = using AI for coding artifacts ethically + effectively (good collaboration); limitation = assumes individual designers are still in power, and takes vibe coding as given rather than critiquing the practice itself.")

n = save("/Users/haukesandhaus/Documents/GitHub/Vibe-Coding-Class/website/slides/Intro-Good-Code-Good-Vibes.pptx")
print("Intro saved ·", n, "slides")
