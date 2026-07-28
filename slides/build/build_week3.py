#!/usr/bin/env python3
"""Week 3 — one deck. Mon 7/27 usability, Tue 7/28 VAP values (RICHER: Discovery +
3-question check + Implementation conflict-resolution, worked example), Wed 7/29 working
day, Thu 7/30 final presentations. VAP content adapted from the speed-run deck + syllabus."""

from deck_common import *
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
prs = init_deck()

FOOT = "Good Code, Good Vibes · TECHIE 1121 · Cornell Tech · Summer 2026"

def numbered_two_col(kick, title, items, foot=None):
    """items: list of (label, gloss). Split into two columns."""
    s = slide(); kicker(s, kick); accent(s)
    text(s, Inches(0.7), Inches(1.3), Inches(12), Inches(1.0),
         [[(title, {"size": 32, "bold": True, "color": WHITE})]])
    half = (len(items) + 1) // 2
    cols = [items[:half], items[half:]]
    xs = [Inches(0.72), Inches(6.9)]
    for ci, col in enumerate(cols):
        paras = []
        for gi, (label, gloss) in enumerate(col):
            n = ci * half + gi + 1
            paras.append([("%02d  " % n, {"color": GREEN, "bold": True, "font": F_MONO, "size": 15}),
                          (label + "  ", {"color": WHITE, "bold": True, "size": 16}),
                          (gloss, {"color": MUTED, "size": 14})])
        text(s, xs[ci], Inches(2.7), Inches(5.9), Inches(4.0), paras, spacing=1.35)
    if foot: footer(s, foot)
    return s

# ===========================================================================
# MONDAY 7/27 - Usability Testing: Human & Agent (video deliverable)
# ===========================================================================
divider("WEEK 3 - MONDAY", "Usability Testing: Can People Use It?",
        "Human + agent testing on your Project 2 prototype")

content("THE TURN", "You built it. Does it work - for someone who isn't you?",
    [[("Weeks 1-2 you built. This week you evaluate, then build the final. Two testing days on ", {"size": 21, "color": WHITE}),
      ("last week's Project 2 prototype:", {"size": 21, "color": WHITE, "bold": True})],
     [("", {})],
     [("Today - ", {"size": 21}), ("can people use it?", {"color": GREEN, "bold": True, "size": 21}),
      ("  (usability).    Tomorrow - ", {"size": 21}),
      ("does it do the right thing?", {"color": GREEN, "bold": True, "size": 21}), ("  (values).", {"size": 21})]], FOOT)

content("WHY TEST", "You are not your user",
    [[("You know where every button is and what every label means - so you are the worst person alive to judge whether your own app is usable.", {"size": 22, "color": WHITE})],
     [("", {})],
     [("Watching ", {"size": 20, "color": MUTED}), ("one", {"size": 20, "color": GREEN, "bold": True}),
      (" real person try it tells you more than a week of guessing. You test the design, not the person - confusion is the design's fault, and it is your best data.", {"size": 20, "color": MUTED})]], FOOT)

big_question("WHAT IT IS", [
    "Give one real person one real task.",
    "Ask them to think aloud.",
    "Stay quiet and watch what they DO."], FOOT)

s = content("WATCH - A REAL USABILITY TEST", "See it done before you do it",
    [[("A moderated think-aloud test. User goal: ", {"size": 20, "color": WHITE}),
      ("order medications online", {"size": 20, "color": GREEN, "bold": True}),
      (" - refill Vitamin D and Tylenol, check which card is on file, set up auto-refill.", {"size": 20, "color": WHITE})],
     [("", {})],
     [("Watch for: where they hesitate, what they say they expect, and how the moderator stays neutral and lets them struggle.", {"size": 18, "color": MUTED})],
     [("", {})],
     [("|>  youtube.com/watch?v=EH7Fx9rpC0c", {"size": 22, "color": GREEN, "bold": True, "font": F_MONO})]], FOOT)
text(s, Inches(0.7), Inches(6.6), Inches(12), Inches(0.32),
     [[("Demo video + usability-study framing adapted from S. Azenkot, INFO 5305 UX & User Research (Cornell Tech)", {})]],
     size=9.5, color=DIM, font=F_MONO)
notes(s, "Play EH7Fx9rpC0c (a few minutes). Pause to point out the think-aloud, the neutral moderation, and one breakdown. Source: Shiri Azenkot, 'Running a Usability Study'.")

numbered_two_col("READING - NIELSEN", "Nielsen's 10 usability heuristics",
    [("Visibility of system status", "always show what's happening"),
     ("Match the real world", "speak the user's language"),
     ("User control & freedom", "easy undo / exits"),
     ("Consistency & standards", "follow conventions"),
     ("Error prevention", "stop mistakes before they happen"),
     ("Recognition over recall", "show options, don't make them remember"),
     ("Flexibility & efficiency", "shortcuts for experts"),
     ("Aesthetic & minimalist", "no needless clutter"),
     ("Recover from errors", "plain-language fixes"),
     ("Help & documentation", "findable when needed")], FOOT)
notes(prs.slides[-1], "Reading: Nielsen - 10 Usability Heuristics (NN/g). Use these as the lens for naming what broke.")

bullets("HOW TO RUN IT", "A usability test in 4 moves",
    ["Give a real task (\"sign up and log your first entry\") - not a tour.",
     "Watch, don't help. Let them struggle a bit - silence is data.",
     "Ask them to think aloud - narrate what they see, expect, and want.",
     "Note every hesitation, wrong turn, and dead end - name it against Nielsen's heuristics."], FOOT)

s = columns("MODERATION", "Prompt without leading", [
    ("ENCOURAGE", "\"What are you thinking?\"", "Keep them talking without steering. Let silence sit, then nudge gently."),
    ("EXPECTATION", "\"What did you expect to happen?\"", "Surfaces the gap between their mental model and your design."),
    ("NEVER", "\"Did you like it?\"", "Leading and unstructured. Watch what they DO - don't fish for approval.")], FOOT)
text(s, Inches(0.7), Inches(6.6), Inches(12), Inches(0.32),
     [[("Moderation guidelines adapted from S. Azenkot, INFO 5305 (Cornell Tech)", {})]],
     size=9.5, color=DIM, font=F_MONO)

s = content("WHAT TO LOOK FOR", "Capture evidence, not your interpretation",
    [[("Look for usability problems - and write down what happened, not what you think it means:", {"size": 20, "color": WHITE})],
     [("", {})],
     [("Breakdowns", {"size": 19, "color": GREEN, "bold": True}), (" - they get stuck, backtrack, or ask for help.", {"size": 19, "color": WHITE})],
     [("Errors", {"size": 19, "color": GREEN, "bold": True}), (" - wrong action, wrong turn, wrong mental model.", {"size": 19, "color": WHITE})],
     [("Workarounds", {"size": 19, "color": GREEN, "bold": True}), (" - they reach the goal, but not the way you designed.", {"size": 19, "color": WHITE})],
     [("", {})],
     [("Note the time, what they did, and a short quote. Refine from the recording afterwards.", {"size": 17, "color": MUTED})]], FOOT)
text(s, Inches(0.7), Inches(6.6), Inches(12), Inches(0.32),
     [[("Evidence framing adapted from S. Azenkot, INFO 5305 (Cornell Tech)", {})]],
     size=9.5, color=DIM, font=F_MONO)

section("PART B", "Now let an AI agent test it", color=GREEN,
        sub="An agent takes browser control of your live site and reports what breaks")

columns("PART B - AGENT TESTING", "An AI drives your live site", [
    ("WHAT", "Agent-based usability test", "Give an AI agent your live URL and the same task. It clicks through and reports where it gets stuck."),
    ("WHY", "Fast, tireless, literal", "Catches broken flows, dead links, unclear labels, impossible steps - in minutes, at any hour."),
    ("TOOLS", "Pick one you can run", "Antigravity | Chrome DevTools MCP | Playwright MCP | browser-use. Log your prompts in history.md.")], FOOT)

columns("HUMAN vs AGENT", "Two testers, different eyes", [
    ("HUMANS CATCH", "Confusion & feeling", "Hesitation, frustration, delight, misread labels, real-world mismatch, the value bending. Emotion you can't fake."),
    ("AGENTS CATCH", "Breakage & coverage", "Dead ends, broken buttons, missing states, every link - fast and repeatable. No fatigue, no mercy."),
    ("AGENTS MISS", "Being human", "They don't feel stuck, don't get embarrassed, don't bring context. Speed is not understanding.")], FOOT)

checklist("THE DELIVERABLE - ONE VIDEO", "Compress hours of testing into 5 minutes", [
    "Record >= 20 min of real testing (2 human tests + 1 agent test).",
    "Cut and speed up to <= 5 minutes - keep the moments that show a finding.",
    "Subtitles required (YouTube auto-captions are fine).",
    "End with findings slides + your spoken summary.",
    "Host unlisted on YouTube; link it in vibe-report.md. No coding today."], FOOT)

s = discuss("ACTIVITY - TODAY", "Test your team's Project 2 prototype - on your own.",
    ["Part A: 2 real testers (not teammates), think-aloud, ~10 min each, screen-recorded.",
     "Part B: 1 agent takes browser control of the live site; log your prompts.",
     "Edit it into one <=5-min subtitled video, then write up human vs agent in vibe-report.md."],
    "Individual - on the group's Project 2 - deliverable = one video + one write-up (due Mon 7/27)", FOOT)
notes(s, "Individual work on the shared Project 2 prototype. Everyone submits their own video + write-up.")

refs("SOURCES", "Watch, test, cite", [
    ("Nielsen - 10 Usability Heuristics (NN/g)", "nngroup.com/articles/ten-usability-heuristics/"),
    ("UEQ - User Experience Questionnaire", "ueq-online.org"),
    ("Usability test demo video", "youtube.com/watch?v=EH7Fx9rpC0c"),
    ("Agent browser control - browser-use", "github.com/browser-use/browser-use"),
    ("Chrome DevTools MCP", "github.com/ChromeDevTools/chrome-devtools-mcp"),
    ("Usability-study material adapted from S. Azenkot, INFO 5305 UX & User Research", "Cornell Tech")], FOOT)


bullets("TOOLS FOR TODAY", "Record, drive, edit, caption",
    ["Record: QuickTime (Mac) / Xbox Game Bar (Win) / Loom / OBS - screen + mic together.",
     "Agent testing: Antigravity, Chrome DevTools MCP, Playwright MCP, or browser-use - an AI drives your live URL.",
     "Edit + speed up: CapCut, iMovie, or Clipchamp - cut to <=5 min, speed the slow parts.",
     "Captions: upload to YouTube (unlisted), turn on auto-captions.",
     "Real UX tools (optional): Maze or Lyssna for prototype tests, Lookback for moderated sessions."], FOOT)
media("SCREENSHOT SLOT", "Show a real test in action",
    "[ screenshot / clip still goes here ]",
    "e.g. a tester mid-task, or your AI agent driving the live site", FOOT)

save("/Users/haukesandhaus/Documents/GitHub/Vibe-Coding-Class/website/slides/Week3-Mon-User-Testing-I.pptx"); prs = init_deck()
# ===========================================================================
# TUESDAY 7/28 - Value Verification: does the design actually change anything?
# ===========================================================================
# ===========================================================================
# TUESDAY 7/28 - Values at Play: the whole framework (3h, with breaks)
# ===========================================================================
divider("WEEK 3 - TUESDAY", "Values at Play: The Whole Framework",
        "Discovery -> Implementation -> Verification -> then verify your own")

content("TODAY", "From \"can people use it?\" to \"what values did we build in - and did they land?\"",
    [[("Monday asked whether people can use your Project 2. Today we go all the way around the Values at Play loop - the method behind everything you have been doing - and then you verify your own value.", {"size": 21, "color": WHITE})],
     [("", {})],
     [("Three parts, with breaks: Discovery (what values?) -> Implementation (build them in) -> Verification (did they land?). Then the canvas, real projects, and your activity.", {"size": 18, "color": MUTED})]], FOOT)

# ---------- FINAL PROJECT KICKOFF - join your team ----------
content("FINAL PROJECT STARTS TODAY", "Commit to your value - then join your team",
    [[("Your team name IS your value. We suggest carrying your Project 2 value through the final - but you can switch, or add a secondary value.", {"size": 21, "color": WHITE})],
     [("", {})],
     [("One person accepts on GitHub Classroom and creates the team with that exact name; everyone else JOINS it - no duplicates.", {"size": 19, "color": MUTED})]], FOOT)

numbered_two_col("JOIN YOUR TEAM", "Six teams - named for the value you carry",
    [("Security", "Isa . Om . Jason"),
     ("Safety & Autonomy", "John . Ajin"),
     ("Sustainability & Trust", "Emily . Kylie . Aria"),
     ("Care & Wellbeing", "Elaine . Winnie . Vienna"),
     ("Productivity", "Liam . Justin . Sebastien"),
     ("Sustainability & Transparency", "Evan . Derin . Magnes . Oliver")], FOOT)
notes(prs.slides[-1], "Two 'Sustainability' teams - check the second word (Trust vs Transparency). Accept + join: classroom.github.com/a/VTWqWkXg. Final Project page (per-team showcase): cornell-tech-vibe-coding-summer-2026.github.io/cornell-tech-vibe-coding-summer-2026-final-project-template/")

# ---------- PART 1 - FOUNDATIONS ----------
section("PART 1 - FOUNDATIONS", "Do artifacts have politics?", sub="Before you design values IN, see that they are already there")

big_question("WINNER, 1980", ["Do artifacts", "have politics?"], FOOT)

media("CASE - ROBERT MOSES", "Overpasses that sorted who reached the beach",
    "[ Robert Moses / Long Island parkway overpass - add image ]",
    "The low overpasses kept buses - and the people who rode them - off the parkway. A design choice, carrying a politics.", FOOT)

content("VALUES LIVE IN THE THING", "A $2.90 rule, cast in steel",
    [[("A subway turnstile does not argue about fare policy - it enforces it. The value (pay to ride) is built into the hardware.", {"size": 21, "color": WHITE})],
     [("", {})],
     [("Winner: an artifact can settle a political question, or require a political arrangement just to work at all.", {"size": 17, "color": MUTED})]], FOOT)

media("CASE - FACEBOOK ALT-TEXT", "\"Image may contain: 1 person, sky, crowd\"",
    "[ Facebook automatic alt-text example - add image ]",
    "What should an AI label - gender? race? Facebook set an 80% confidence bar, higher for sensitive tags. Every one of those is a value choice.", FOOT)

photo_split("THE PRACTICAL TURN", "The conscientious designer",
    [[("Once values are always in the thing, the question flips: not \"are there values?\" but \"which ones - and did I choose them on purpose?\"", {"size": 19, "color": WHITE})],
     [("", {})],
     [("Values at Play is the method for doing that deliberately.", {"size": 18, "color": GREEN, "bold": True})]],
    "vap_book.png", "Flanagan & Nissenbaum, Values at Play (MIT Press) - VAP course, Engelmann & Nissenbaum, Cornell Tech", FOOT)

section("THE HEURISTIC", "Discovery -> Implementation -> Verification", sub="One loop, three moves - and it is cyclical")

content("BREAK", "~10 minutes",
    [[("Stretch, water. We pick up with Discovery: what values are even in play?", {"size": 22, "color": WHITE})]], FOOT)

# ---------- PART 2 - DISCOVERY ----------
section("PART 2 - DISCOVERY", "What values? And what do they actually mean?")

columns("DISCOVERY I - CONSIDER THE SOURCES", "Where the values hide", [
    ("FUNCTIONAL DESCRIPTION", "What does it do?", "The core function already implies values - a secure messenger implies privacy; a feed implies attention."),
    ("KEY ACTORS", "Who is involved?", "Designers, funders, users, bystanders. Whose interests are in the room - and whose are not?"),
    ("CONTEXT + CONSTRAINTS", "Where does it live?", "Societal norms, laws, and technical limits (screen size, bandwidth, model context) all carry values.")], FOOT)
notes(prs.slides[-1], "Discovery I heuristic (VAP, Engelmann & Nissenbaum). The upshot is a LIST OF VALUES - including 'collateral' values you didn't set out to touch.")

content("DISCOVERY II - OPERATIONALIZE", "A value is not one thing",
    [[("An abstract value has to be defined before you can build it or test it:", {"size": 20, "color": WHITE})],
     [("", {})],
     [("Privacy", {"size": 19, "color": GREEN, "bold": True}), (" = control?  secrecy?  contextual integrity?", {"size": 19, "color": WHITE})],
     [("Fairness", {"size": 19, "color": GREEN, "bold": True}), (" = equal shares?  or proportional to merit / need?", {"size": 19, "color": WHITE})],
     [("", {})],
     [("Your operational definition decides your design AND your test. (That is yesterday's move - and today's.)", {"size": 17, "color": MUTED})]], FOOT)

columns("CASE - WHAT'S FAIR?", "One value, incompatible definitions", [
    ("THE SETUP", "Split a limited loan", "Person A repays 100%, Person B repays 20%. Identical otherwise. How do you split the money?"),
    ("THE OPTIONS", "All to A? Or share?", "All to A maximizes repayment - and leaves B nothing. Proportional shares the pot but funds a worse bet."),
    ("THE POINT", "No free lunch", "A single algorithm can't satisfy every definition of fairness. Someone chooses - make it conscious.")], FOOT)

content("BREAK", "~10 minutes",
    [[("Back soon - next we turn values into features, and handle the clashes.", {"size": 22, "color": WHITE})]], FOOT)

# ---------- PART 3 - IMPLEMENTATION ----------
section("PART 3 - IMPLEMENTATION", "Translate values into features - and resolve the clashes")

photo_split("TRANSLATE", "From a value to a spec",
    [[("Implementation turns an ethical or political value into concrete requirements, features, and architecture.", {"size": 19, "color": WHITE})],
     [("", {})],
     [("Mary Flanagan's Giant Joystick makes \"play\" physical - one value, rebuilt into the hardware itself.", {"size": 17, "color": MUTED})]],
    "giant_joystick.png", "Mary Flanagan, 'Giant Joystick' - Values at Play (educational use)", FOOT)

content("DESIGN CHOICES CARRY VALUES", "The same feed, two values",
    [[("Design a social-media feed for ", {"size": 20, "color": WHITE}), ("engagement", {"size": 20, "color": RED, "bold": True}),
      (", and you get autoplay, streaks, and notifications - \"push people's buttons.\"", {"size": 20, "color": WHITE})],
     [("Design it for ", {"size": 20, "color": WHITE}), ("autonomy", {"size": 20, "color": GREEN, "bold": True}),
      (", and you add friction, stopping points, and real user control.", {"size": 20, "color": WHITE})],
     [("", {})],
     [("Same product. The values show up in the defaults.", {"size": 17, "color": MUTED})]], FOOT)

columns("WHEN VALUES CLASH", "Name the resolution - out loud", [
    ("DISSOLVE", "Redesign around it", "Find a design where BOTH values hold. (Vacuum: on-device maps, nothing to the cloud.)"),
    ("COMPROMISE", "Each in part", "Promote each value in less than full measure. (Vacuum: cloud data, but anonymized + opt-out.)"),
    ("TRADE-OFF", "Sacrifice one", "Give one up for another - and say so out loud. (Vacuum: full collection for max automation.)")], FOOT)

content("BREAK", "~10 minutes",
    [[("Last stretch. Then: did any of it actually work?", {"size": 22, "color": WHITE})]], FOOT)

# ---------- PART 4 - VERIFICATION ----------
section("PART 4 - VERIFICATION", "Did the values actually land?")

content("VERIFICATION", "Assessing whether the values made it in",
    [[("\"Verification involves assessing whether efforts to integrate values have succeeded.\" Cyclical - ask at every step.", {"size": 21, "color": WHITE})],
     [("", {})],
     [("Three ways to ask - and your operationalization decides which one fits.", {"size": 17, "color": MUTED})]], FOOT)

s = columns("THE THREE LENSES", "Behavior - Understanding - Affect", [
    ("BEHAVIOR", "Do they ACT differently?", "Observation, A/B, before/after, task completion."),
    ("UNDERSTANDING", "Do they GRASP it?", "Survey, interview, think-aloud, scenario."),
    ("AFFECT", "Do they FEEL differently?", "Attitude survey, diary, interview, choice under pressure.")], FOOT)
text(s, Inches(0.7), Inches(6.62), Inches(12), Inches(0.32),
     [[("Verification lenses: Values at Play - VAP course, Engelmann & Nissenbaum (Cornell Tech)", {})]],
     size=9.5, color=DIM, font=F_MONO)

photo_split("EXAMPLE - BEHAVIOR", "Does the design change what people DO?",
    [[("A smart meter only \"works\" for sustainability if people actually use less. Behavior is the sharpest test.", {"size": 18, "color": WHITE})],
     [("", {})],
     [("Verify how: prototype - A/B - before/after - task completion.", {"size": 16, "color": MUTED})]],
    "smart_meter.jpg", "Behavior verification (Values at Play) - smart-meter image, educational use", FOOT)

bullets("PICK THE INSTRUMENT THAT FITS YOUR VALUE", "Examples - tie the tool to the value you are analyzing",
    ["Beauty / craft / originality -> the Vibed Slop Detector (is it generic AI slop?), or UEQ attractiveness.",
     "Autonomy / competence / wellbeing -> METUX + Self-Determination scales.",
     "Trust / warmth / safety -> a validated scale from the HRI Scale Database.",
     "Overall UX beyond usability -> UEQ. The instrument has to match the value - don't default to UEQ."], FOOT)

quote("FROM OUR RESEARCH",
    "Persuasion-focused evaluation nearly DOUBLED how often professional designers rejected dark-pattern designs - and shifted their reasoning from business justifications toward user autonomy and well-being.",
    "Sandhaus, Rhomberg & Nissenbaum, CHIWORK 2026 - a study of 141 UX professionals",
    "Measuring a value doesn't just describe it - it changes what gets built. That is why verification is not optional.", FOOT)

# ---------- PART 5 - THE CANVAS + PROJECTS ----------
section("PART 5 - IN PRACTICE", "The VAP canvas + real projects")

media("THE FIGJAM CANVAS", "One board, the whole loop",
    "[ walk the 'Vibing Values' canvas live ]",
    "figma.com/board/.../Vibing-Values - each tab is one move (Discovery / Implementation / Verification). We walk them as connected case studies.", FOOT)

media("PITECH PROJECTS", "Fellows who ran the whole framework",
    "[ show the PiTech projects + website - add screenshots ]",
    "Real value-driven projects, and what worked - and what was hard - when teams used VAP end to end.", FOOT)

# ---------- PART 6 - BACK TO YOUR ACTIVITY ----------
s = discuss("NOW YOU (TEAM)", "Verify the value your Project 2 claims.",
    ["State how you operationalize your value - that definition decides the method.",
     "No time to run tests: RESEARCH how each lens could be verified (behavior especially) + what design interventions have worked.",
     "Does it deliver the value - for whom, and who's left out? Then: what would you change?"],
    "A value-centered user RESEARCH deck + a reflection on FigJam - each person owns a slide + a lane - due today", FOOT)
notes(s, "Hand back to the activity. Language note: no 'verdict' - just 'does it deliver the value, and for whom?'")

content("FIND YOUR PAPER - EACH PERSON", "Section 04 is individual: real evidence, found with AI search",
    [[("For your named slide, find ONE real paper on a design intervention that tried to move your value:", {"size": 20, "color": WHITE})],
     [("", {})],
     [("Claude / ChatGPT (ask for real, citable papers - then open + check them), Elicit, Semantic Scholar, Google Scholar - and its new Scholar Labs.", {"size": 19, "color": GREEN, "bold": True})],
     [("", {})],
     [("Never invent a citation. Open the actual paper and confirm it says what you think.", {"size": 17, "color": MUTED})]], FOOT)

section("QUIET READING", "15 minutes - read the paper YOU found", color=GREEN,
        sub="Phones away. Individual. This is the discovery reading - it fills your section-04 slide (FigJam: 04 Outlook).")

content("WHILE YOU READ", "Pull out what your slide needs",
    [[("From the paper you found, mark:", {"size": 20, "color": WHITE})],
     [("1.  the design intervention they tried - what exactly did they change?", {"size": 18, "color": MUTED})],
     [("2.  the result - did it move the value? worked / mixed / failed", {"size": 18, "color": MUTED})],
     [("3.  whether it supports OUR design direction", {"size": 18, "color": MUTED})],
     [("", {})],
     [("Then put it on your named slide (section 04), with the citation.", {"size": 17, "color": DIM})]], FOOT)
notes(prs.slides[-1], "Protected, individual reading of the paper each person found via AI academic search (Claude / ChatGPT / Elicit / Semantic Scholar / Google Scholar Labs). This is the discovery reading the timer is built around - keep it silent.")

refs("SOURCES", "Read + cite", [
    ("Flanagan & Nissenbaum - Values at Play in Digital Games (MIT Press, 2014)", "mitpress.mit.edu"),
    ("Winner - Do Artifacts Have Politics? (Daedalus, 1980)", "the Robert Moses + turnstile cases"),
    ("van de Poel - Translating Values into Design Requirements (2013)", "implementation"),
    ("VAP course - S. Engelmann & H. Nissenbaum", "Cornell Tech - INFO 5010 (PiTech Ethics)"),
    ("Sandhaus, Rhomberg & Nissenbaum 2026 - User Testing Promotes Ethical Sensibility (CHIWORK)", "osf.io/nw2tj"),
    ("UEQ - ueq-online.org  |  HRI Scale DB  |  METUX / SDT scales", "verification instruments"),
    ("Vibed Slop Detector + impeccable.style slop rules", "github.com/HaukeCornell/Vibed-Slop-Detector")], FOOT)

save("/Users/haukesandhaus/Documents/GitHub/Vibe-Coding-Class/website/slides/Week3-Tue-User-Testing-II.pptx"); prs = init_deck()
# ===========================================================================
# WEDNESDAY 7/29 - Planning Day (Final Project, Day 2): What's the problem?
# ===========================================================================
divider("WEEK 3 - WEDNESDAY", "Planning Day: What's the problem?",
        "Final project - value-centered need finding + storyboard")

content("THE TURN", "Tuesday told you IF it works. Today you decide what to build.",
    [[("You verified your Project 2 and found where the design doesn't move the value. Today you plan the fix - for a real person ", {"size": 21, "color": WHITE}),
      ("outside this class.", {"size": 21, "color": GREEN, "bold": True})],
     [("", {})],
     [("Open with Remy Stewart's Figma guest talk, then find the real need.", {"size": 19, "color": MUTED})]], FOOT)

content("KEEP YOUR VALUE", "Change the design, not the commitment",
    [[("You already have a value and what Tuesday's verification showed. The final is the NEXT iteration of that value - not a new theme.", {"size": 21, "color": WHITE})],
     [("", {})],
     [("Today's job: name a specific beneficiary, and find what they actually need.", {"size": 19, "color": MUTED})]], FOOT)

columns("FIND THE NEED", "Talk to someone - don't guess", [
    ("INTERVIEW", "Ask, don't pitch", "A short conversation with a real person outside the class. Ask about their problem, not your idea."),
    ("PERSONA", "One specific human", "Name, situation, goal, frustration, context. Not 'students' - one person you design for."),
    ("THE NEED", "In their words", "Write the need the way they would say it. That sentence is what your design has to serve.")], FOOT)

content("PERSONAS - WHY, AND THE RISK", "A persona is a tool to EMPATHIZE - not to stereotype",
    [[("A persona holds one real user's needs in mind while you build - so you design for someone specific, not \"everyone.\"", {"size": 20, "color": WHITE})],
     [("", {})],
     [("The risk: personas flatten people into cliches. Used well they build empathy; used badly they just launder your own assumptions.", {"size": 18, "color": MUTED})]], FOOT)

content("GET AS CLOSE TO THE PERSON AS YOU CAN", "Nothing beats real contact",
    [[("Best: talk to a real person - a short interview. Better still: co-design, where the user helps make the thing.", {"size": 20, "color": WHITE})],
     [("", {})],
     [("We can't fully do participatory design here - so get as close as you can, and be honest about the gap.", {"size": 18, "color": MUTED})]], FOOT)

s = discuss("AI PERSONAS - THEATER OR INSIGHT?", "Make an AI persona, then hold it against a real human.",
    ["Generate a persona with AI - even \"chat\" with the simulated user.",
     "Compare it to the real person you interviewed: what did each get right?",
     "Does the AI theater help you empathize - or hand your own stereotypes back as if they were data?"],
    "Bring BOTH into planning: the AI persona AND one real human's words", FOOT)
notes(s, "AI personas are fast + frictionless but can launder assumptions back as 'data'. Comparing to a real human is the check. Ref: nngroup.com/articles/synthetic-users/")

media("SCREENSHOT SLOT", "AI persona vs. real interview",
    "[ AI-generated persona + simulator chat, next to your real interview notes ]",
    "add the two side by side - where did the AI flatten the person?  (skeptical take: NN/g - Synthetic Users)", FOOT)

content("KEEP THE VALUE IN THE FOREGROUND", "Values-conscientious planning",
    [[("Put your value at the front of BOTH artifacts:", {"size": 20, "color": WHITE})],
     [("", {})],
     [("Persona", {"size": 19, "color": GREEN, "bold": True}), (" - where does this user MEET or LOSE the value in their day?", {"size": 19, "color": WHITE})],
     [("Storyboard", {"size": 19, "color": GREEN, "bold": True}), (" - show the value living in each frame: the moment it's supported, defended, or repaired.", {"size": 19, "color": WHITE})]], FOOT)

content("OPERATIONALIZE - THE VALUE TREE", "From value down to the prompts you'll build",
    [[("Break your value down until it becomes something you can actually build:", {"size": 20, "color": WHITE})],
     [("", {})],
     [("value  ->  norms  ->  requirements  ->  features  ->  ", {"size": 22, "color": WHITE}),
      ("the prompts you run tomorrow", {"size": 22, "color": GREEN, "bold": True})],
     [("", {})],
     [("The bottom of the tree IS your Thursday build plan - the exact prompts, grounded in the value.", {"size": 17, "color": MUTED})]], FOOT)

s = bullets("STORYBOARD IT", "4-8 frames: the value in use",
    ["Frame the user, their problem, where your thing enters, and what changes.",
     "Hand-drawn (photograph it) or AI-assisted - rough is fine.",
     "AI can sketch or fill frames fast - but keep YOUR story; the AI is a co-designer, not the author.",
     "It doubles as your Thursday pitch skeleton - show the value living in the design, not just the copy."], FOOT)
notes(s, "AI storyboarding / co-design: Sandhaus, Gu, Parreira & Ju 2025, 'Co-Designing with Transformers' (DIS '25) - GenAI helps generate + iterate frames, but can flatten the narrative if you let it drive. dl.acm.org/doi/10.1145/3715336.3735805")

bullets("TOOLS FOR TODAY", "Storyboard + pitch",
    ["Storyboard: FigJam / Figma, Canva, Google Slides (one frame per slide), or Storyboarder (free & open-source).",
     "AI pitch video (for Thursday): Google Veo 3.1 (in Flow), Runway, Kling, or Pika.",
     "Note: OpenAI retired the Sora app in April 2026 - use Veo / Runway / Kling / Pika instead.",
     "Free tiers add watermarks + limits - fine for a class pitch.",
     "Interview + persona: just talk to a person; capture notes in FigJam or a doc."], FOOT)
notes(prs.slides[-1], "Sora web/app discontinued 2026-04-26. Recommend Veo 3.1 / Runway / Kling / Pika for pitch videos.")

media("SCREENSHOT SLOT", "Show an example",
    "[ example storyboard or persona goes here ]",
    "a strong persona, a storyboard, or an AI-pitch still", FOOT)

discuss("ACTIVITY - TODAY (TEAM)", "Plan the next iteration of your value.",
    ["Interview a real person outside the class -> build a persona + the need in their words.",
     "Operationalize: value -> norms -> requirements -> features -> the prompts you'll build.",
     "Storyboard the solution (4-8 frames) and scope the smallest real thing for Thursday."],
    "Deliverable = planning-report.md + your storyboard - this is Checkpoint 2 of the final", FOOT)

content("LAST 30 MINUTES - START VIBE CODING", "Don't wait for Thursday to open the repo",
    [[("Once the plan is locked, spend the final half hour getting a head start:", {"size": 21, "color": WHITE})],
     [("", {})],
     [("Open your team repo, run the FIRST prompt from your value tree, and get anything at all on screen.", {"size": 19, "color": GREEN, "bold": True})],
     [("", {})],
     [("Thursday has no lecture - it's all build - so every minute you start now is a minute you keep tomorrow.", {"size": 17, "color": MUTED})]], FOOT)

content("CARRY IT TO THURSDAY", "Your plan is your build",
    [[("Tomorrow is a full build day (no lecture): build straight from the prompts your value tree produced, then showcase it in the afternoon (3-5).", {"size": 21, "color": WHITE})],
     [("", {})],
     [("Guest lecture today: Remy Stewart (Figma). Bring your questions about design + prototyping.", {"size": 17, "color": MUTED})]], FOOT)

save("/Users/haukesandhaus/Documents/GitHub/Vibe-Coding-Class/website/slides/Week3-Wed-Working-Day.pptx"); prs = init_deck()
# ===========================================================================
# THURSDAY 7/30 - Build day + 3-5 Showcase (Final Project, Day 3+4): What's the solution?
# ===========================================================================
divider("WEEK 3 - THURSDAY", "Build day: What's the solution?",
        "No lecture - build all morning, showcase at 3-5")

content("THE DAY", "No lecture. Heads-down build, then showcase.",
    [[("You already tested Monday (usability) and researched the value Tuesday - today is NOT a new test day. Build the smallest real thing, and tell its story well.", {"size": 21, "color": WHITE})],
     [("", {})],
     [("Everything you planned Wednesday becomes an artifact people can use. Showcase runs 3-5.", {"size": 19, "color": MUTED})]], FOOT)

bullets("BUILD", "The smallest real thing - done well",
    ["Build from the prompts your value tree produced Wednesday.",
     "One complete action a stranger can finish - no broken buttons, no fake features.",
     "The value lives in the design, not just the copy.",
     "Polish the happy path and get it hosted - a stranger will try it at the showcase."], FOOT)

bullets("THE PITCH VIDEO", "Show the IDEA of your app - authentically",
    ["Make a short video that shows the idea of your app.",
     "It does NOT have to be AI: a screen recording, an acted skit, or a hand-animated walkthrough all work.",
     "AI video is fine too: Google Veo 3.1 (Flow), Runway, Kling, Pika. (OpenAI retired the Sora app in 2026.)",
     "Pick whatever feels AUTHENTIC to your team. Put it in presentation/ AND embed it in your final slides."], FOOT)

media("SCREENSHOT SLOT", "Show a build / pitch still",
    "[ demo screenshot or pitch still goes here ]",
    "the app in use, or a frame from the pitch video", FOOT)

bullets("HOW TO PRESENT", "Tell the arc of your value - not a test report",
    ["The artifact - a live demo of the happy path, value living in the design.",
     "The value + who - what it operationalizes and for whom, with your PERSONA and STORYBOARD.",
     "The journey - your research findings, the design ALTERNATIVES you considered, how the design PROGRESSED from Project 2.",
     "The pitch - your embedded video.",
     "The afterlife - who keeps using or maintaining it."], FOOT)

columns("HOW IT'S GRADED", "Four checkpoints, one grade", [
    ("TUE - VERIFY", "Checkpoint 1", "Value verification: research deck + canvas."),
    ("WED - PLAN", "Checkpoint 2", "Value-centered planning: planning report + storyboard."),
    ("THU - BUILD", "Checkpoint 3", "Hosted app + project report + pitch.")], FOOT)

columns("HOW IT'S GRADED", "...and the showcase", [
    ("THU 3-5 - SHOWCASE", "Checkpoint 4", "The live demo + slides at the afternoon showcase."),
    ("PROJECT", "25%", "Summary of the four Canvas checkpoints, against the project rubric."),
    ("PRESENTATION", "15%", "A separate grade for the showcase. Peer-check adjusts for unfair contribution.")], FOOT)
notes(prs.slides[-1], "Grade = summary of the 4 Canvas checkpoints (25% project) + a separate 15% presentation grade. Ethical + individual reflections must be human-written.")

content("LOGISTICS", "Afternoon showcase (3-5)",
    [[("A shared showcase - you present alongside other classes' work, with a keynote for families. Prototype + slides on the monitors; demo from the HOSTED link, not localhost.", {"size": 21, "color": WHITE})],
     [("", {})],
     [("Keep a screen-recording backup in presentation/ in case the Wi-Fi drops. Final tweaks + submission can come after the showcase.", {"size": 18, "color": MUTED})]], FOOT)

content("THAT'S A WRAP", "Good code, good vibes",
    [[("You learned to build fast with AI - and to build right: name the value, test with real people, and own the calls you made.", {"size": 22, "color": WHITE})],
     [("", {})],
     [("Now go make something the world keeps benefiting from.", {"size": 20, "color": GREEN, "bold": True})]], FOOT)

n = save("/Users/haukesandhaus/Documents/GitHub/Vibe-Coding-Class/website/slides/Week3-Thu-Final-Presentations.pptx")
print("Week 3 saved as 4 per-day decks (Thu:", n, "slides)")
